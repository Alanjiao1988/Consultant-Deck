from pathlib import Path
import json
import sys

import yaml
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import add_textbox
from scripts.qa_pptx import load_facts, main as qa_main, numeric_facts, run_qa


def make_deck(path: Path, body: str = "Savings reach 30% in the base case.") -> Path:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Evidence loader test", "Synthetic", "Test", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Automation can reduce target-process cost by 30%")
    add_textbox(slide, body, 1.2, 3.2, 20, 1)
    add_footer(slide, 2, "Synthetic test source")
    prs.save(path)
    return path


def write_briefs(path: Path, minimum: int = 1) -> Path:
    path.write_text(
        yaml.safe_dump(
            {
                "content_density_target": "research-heavy",
                "pages": [
                    {"page": 1, "page_role": "cover"},
                    {
                        "page": 2,
                        "page_role": "core_argument",
                        "min_registered_numbers": minimum,
                    },
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )
    return path


def test_load_facts_preserves_qualitative_records_without_value(tmp_path):
    evidence = tmp_path / "evidence.json"
    evidence.write_text(
        json.dumps(
            {
                "facts": [
                    {
                        "id": "Q001",
                        "claim": "Regulatory approval is required before production use",
                        "source_type": "regulator",
                        "used_on_pages": [2],
                    }
                ],
                "calculations": [],
            }
        ),
        encoding="utf-8",
    )
    records = load_facts(evidence)
    assert len(records) == 1
    assert records[0]["evidence_type"] == "qualitative"
    assert records[0]["is_numeric"] is False
    assert records[0]["value"] is None
    assert records[0]["text_value"].startswith("Regulatory approval")
    assert numeric_facts(records) == []


def test_malformed_records_and_pages_do_not_crash(tmp_path):
    evidence = tmp_path / "evidence.json"
    evidence.write_text(
        json.dumps(
            {
                "facts": [
                    None,
                    "bad-row",
                    {
                        "id": "Q002",
                        "claim": "Named-vendor conclusion",
                        "value": "Vendor A",
                        "used_on_pages": ["2", "bad", None, -1],
                    },
                    {
                        "id": "Q003",
                        "claim": "Boolean is not a numeric fact",
                        "value": True,
                        "used_on_pages": "2",
                    },
                ],
                "calculations": {"not": "a-list"},
            }
        ),
        encoding="utf-8",
    )
    records = load_facts(evidence)
    assert len(records) == 2
    assert all(record["evidence_type"] == "qualitative" for record in records)
    assert records[0]["used_on_pages"] == {2}
    assert records[1]["used_on_pages"] == {2}


def test_qualitative_only_evidence_does_not_satisfy_numeric_density(tmp_path):
    deck = make_deck(tmp_path / "qualitative_only.pptx", "The design requires regulatory approval.")
    evidence = tmp_path / "evidence.json"
    evidence.write_text(
        json.dumps(
            {
                "facts": [
                    {"id": "Q001", "claim": "Regulatory approval is required", "used_on_pages": [2]}
                ],
                "calculations": [],
            }
        ),
        encoding="utf-8",
    )
    briefs = write_briefs(tmp_path / "briefs.yaml")
    findings = run_qa(deck, evidence, briefs)
    assert any(
        item.severity == "error"
        and item.check == "data_density"
        and "No registered numeric facts" in item.message
        for item in findings
    )


def test_mixed_numeric_and_qualitative_evidence_preserves_numeric_qa(tmp_path):
    deck = make_deck(tmp_path / "mixed.pptx")
    evidence = tmp_path / "evidence.json"
    evidence.write_text(
        json.dumps(
            {
                "facts": [
                    {"id": "Q001", "claim": "Regulatory approval is required", "used_on_pages": [2]},
                    {
                        "id": "N001",
                        "claim": "Base-case savings",
                        "value": "30",
                        "unit": "%",
                        "used_on_pages": [2],
                    },
                ],
                "calculations": [],
            }
        ),
        encoding="utf-8",
    )
    briefs = write_briefs(tmp_path / "briefs.yaml")
    findings = run_qa(deck, evidence, briefs)
    assert not any(item.check == "fact_consistency" and item.severity == "error" for item in findings)
    assert not any(item.check == "data_density" for item in findings)


def test_cli_diagnostics_separate_numeric_and_qualitative_counts(tmp_path, capsys):
    deck = make_deck(tmp_path / "diagnostics.pptx")
    evidence = tmp_path / "evidence.json"
    evidence.write_text(
        json.dumps(
            {
                "facts": [
                    {"id": "Q001", "claim": "Qualitative conclusion"},
                    {"id": "N001", "claim": "Savings", "value": 30, "unit": "%", "used_on_pages": [2]},
                ],
                "calculations": [],
            }
        ),
        encoding="utf-8",
    )
    assert qa_main([str(deck), "--facts", str(evidence)]) == 0
    output = capsys.readouterr().out
    assert "facts=2 (numeric=1, qualitative=1)" in output
