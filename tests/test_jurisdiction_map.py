from __future__ import annotations

from io import BytesIO
import json
from pathlib import Path
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.util import Cm
import pytest

ROOT = Path(__file__).resolve().parents[1]

from scripts.demo_generate_jurisdiction_map import generate
from scripts.jurisdiction_map import SVG_EXTENSION_URI, categorical_jurisdiction_map


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _categories():
    return [
        {"id": "a", "label": "Pattern A", "description": "First pattern"},
        {"id": "b", "label": "Pattern B", "description": "Second pattern"},
    ]


def _jurisdictions():
    return [
        {"id": "US", "label": "United States", "category_id": "a", "note": "Federal and state overlay", "evidence_ids": ["F1"]},
        {"id": "EU", "label": "European Union", "category_id": "b", "note": "Bloc rulebook", "evidence_ids": ["F2"], "entity_type": "regulatory_bloc", "members": ["EU members"]},
    ]


def test_demo_embeds_png_fallback_svg_extension_and_native_overlays(tmp_path):
    pptx_path, manifest_path = generate(tmp_path / "map.pptx", tmp_path / "map.exhibits.json")
    with zipfile.ZipFile(pptx_path) as archive:
        png_media = [name for name in archive.namelist() if name.endswith(".png")]
        svg_media = [name for name in archive.namelist() if name.endswith(".svg")]
        assert len(png_media) == 1
        assert len(svg_media) == 1

        png_bytes = archive.read(png_media[0])
        assert png_bytes.startswith(b"\x89PNG\r\n\x1a\n")
        with Image.open(BytesIO(png_bytes)) as image:
            assert image.size == (1000, 520)
        assert len(archive.read(svg_media[0])) > 6_000

        content_types = archive.read("[Content_Types].xml").decode("utf-8")
        assert "image/png" in content_types
        assert "image/svg+xml" in content_types

        slide_xml = archive.read("ppt/slides/slide1.xml").decode("utf-8")
        rels_xml = archive.read("ppt/slides/_rels/slide1.xml.rels").decode("utf-8")
        assert SVG_EXTENSION_URI in slide_xml
        assert "svgBlip" in slide_xml
        assert "drawing/2016/SVG/main" in slide_xml
        assert ".png" in rels_xml
        assert ".svg" in rels_xml

    prs = Presentation(pptx_path)
    names = [shape.name for shape in prs.slides[0].shapes]
    assert sum(name.startswith("MAP_MARKER|") for name in names) == 6
    assert sum(name.startswith("MAP_LABEL|") for name in names) == 6
    assert sum(name.startswith("MAP_LEGEND|") for name in names) == 4

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))["exhibits"][0]
    assert manifest["classification_dimension"]
    assert manifest["coverage"] == {
        "reviewed": 24,
        "shown": 6,
        "selection_basis": "Representative operating patterns",
    }
    assert manifest["base_image"]["primary"] == "png_fallback"
    assert manifest["base_image"]["vector_extension"] == "office_svg_blip"


def test_demo_label_boxes_do_not_materially_overlap(tmp_path):
    pptx_path, _ = generate(tmp_path / "map.pptx", tmp_path / "map.exhibits.json")
    prs = Presentation(pptx_path)
    labels = [shape for shape in prs.slides[0].shapes if shape.name.startswith("MAP_LABEL|")]
    for index, first in enumerate(labels):
        for second in labels[index + 1:]:
            overlap_w = min(first.left + first.width, second.left + second.width) - max(first.left, second.left)
            overlap_h = min(first.top + first.height, second.top + second.height) - max(first.top, second.top)
            if overlap_w <= 0 or overlap_h <= 0:
                continue
            overlap = overlap_w * overlap_h
            smaller = min(first.width * first.height, second.width * second.height)
            assert overlap / smaller < 0.20


def test_unknown_category_is_rejected():
    _, slide = _blank_slide()
    jurisdictions = _jurisdictions()
    jurisdictions[0]["category_id"] = "missing"
    with pytest.raises(ValueError, match="unknown category"):
        categorical_jurisdiction_map(
            slide, jurisdictions, _categories(), 1, 3, 30, 13,
            classification_dimension="Test dimension",
            as_of="2026-07-14",
            caveat="Test caveat",
        )


def test_unused_legend_category_is_rejected():
    _, slide = _blank_slide()
    jurisdictions = [_jurisdictions()[0]]
    with pytest.raises(ValueError, match="legend categories are not referenced"):
        categorical_jurisdiction_map(
            slide, jurisdictions, _categories(), 1, 3, 30, 13,
            classification_dimension="Test dimension",
            as_of="2026-07-14",
            caveat="Test caveat",
        )


def test_custom_anchor_requires_explicit_rationale():
    _, slide = _blank_slide()
    jurisdictions = [
        {
            "id": "NEW",
            "label": "New market",
            "category_id": "a",
            "note": "New classification",
            "evidence_ids": ["F1"],
            "anchor": {"x": 0.5, "y": 0.5},
        },
        {
            "id": "EU",
            "label": "European Union",
            "category_id": "b",
            "note": "Bloc rulebook",
            "evidence_ids": ["F2"],
            "entity_type": "regulatory_bloc",
            "members": ["EU members"],
        },
    ]
    with pytest.raises(ValueError, match="custom_anchor_rationale"):
        categorical_jurisdiction_map(
            slide, jurisdictions, _categories(), 1, 3, 30, 13,
            classification_dimension="Test dimension",
            as_of="2026-07-14",
            caveat="Test caveat",
            allow_custom_anchor=True,
        )


def test_world_assets_include_vector_and_raster_fallback():
    svg = (ROOT / "assets" / "maps" / "world_equal_earth.svg").read_text(encoding="utf-8")
    png = ROOT / "assets" / "maps" / "world_equal_earth.png"
    assert 'viewBox="0 0 1000 520"' in svg
    assert svg.count("<path") >= 30
    assert len(svg) > 6_000
    assert png.read_bytes().startswith(b"\x89PNG\r\n\x1a\n")
    with Image.open(png) as image:
        assert image.size == (1000, 520)
