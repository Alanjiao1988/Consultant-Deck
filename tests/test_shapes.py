from pathlib import Path
import json
import sys
import zipfile
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_W, SLIDE_H, add_action_title, add_footer, cover_page
from scripts import consulting_shapes as shapes
from scripts.consulting_shapes import add_textbox, waterfall, native_chart
from scripts.demo_generate_deck import generate as generate_demo
from scripts.qa_pptx import run_qa


def blank_deck():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def save_facts(path, facts):
    path.write_text(json.dumps({"facts": facts}, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def test_font_xml_contains_latin_and_ea(tmp_path):
    prs, slide = blank_deck()
    add_action_title(slide, "分三波迁移可在 18 个月内降低 TCO")
    add_textbox(slide, "中文 English 123", 1, 3, 10, 1)
    out = tmp_path / "font.pptx"
    prs.save(out)
    with zipfile.ZipFile(out) as zf:
        xml = "\n".join(zf.read(name).decode("utf-8", errors="ignore") for name in zf.namelist() if name.startswith("ppt/slides/slide"))
    assert 'typeface="Arial"' in xml
    assert 'typeface="Microsoft YaHei"' in xml


def test_waterfall_math_validation():
    _, slide = blank_deck()
    waterfall(slide, [("Base", 100, "start"), ("Saving", -20, "delta"), ("Target", 80, "end")], 1, 3, 10, 6)


def test_native_chart_validates_series_length():
    _, slide = blank_deck()
    try:
        native_chart(slide, "bar", ["A", "B"], [("x", [1])], 1, 3, 10, 5)
    except ValueError:
        return
    raise AssertionError("Expected ValueError")


def test_theme_json_matches_shape_constants():
    theme = json.loads((ROOT / "assets" / "theme.json").read_text(encoding="utf-8"))
    assert theme["fonts"]["latin"] == shapes.FONT_LATIN
    assert theme["fonts"]["east_asian"] == shapes.FONT_EA
    assert theme["colors"]["primary"].lstrip("#") == shapes.PRIMARY
    assert theme["colors"]["accent"].lstrip("#") == shapes.ACCENT
    assert theme["colors"]["text"].lstrip("#") == shapes.GRAY_TEXT
    assert theme["colors"]["secondary_text"].lstrip("#") == shapes.GRAY_2
    assert theme["colors"]["light_text"].lstrip("#") == shapes.GRAY_3
    assert theme["colors"]["line"].lstrip("#") == shapes.GRAY_LINE
    assert theme["colors"]["fill"].lstrip("#") == shapes.GRAY_FILL


def test_demo_deck_has_no_qa_false_positive_findings_with_facts(tmp_path):
    out = tmp_path / "demo_ai_transformation.pptx"
    facts = tmp_path / "demo_ai_transformation.evidence.json"
    generate_demo(out, facts)
    findings = run_qa(out, facts)
    assert findings == []


def test_mixed_cover_terminology_is_checked(tmp_path):
    prs, slide = blank_deck()
    cover_page(slide, "某银行数字化转型 strategy", "Mixed-language cover", "Client", "2026-07-05")
    out = tmp_path / "mixed_cover.pptx"
    prs.save(out)
    findings = run_qa(out)
    assert any(f.check == "terminology" and f.slide == 1 for f in findings)


def test_pure_english_cover_has_no_terminology_warning(tmp_path):
    prs, slide = blank_deck()
    cover_page(slide, "Digital transformation strategy", "English cover", "Client", "2026-07-05")
    out = tmp_path / "english_cover.pptx"
    prs.save(out)
    findings = run_qa(out)
    assert not any(f.check == "terminology" for f in findings)


def test_registered_fact_passes_consistency_check(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Fact test deck", "English cover", "Client", "2026-07-05")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Savings can reach 30% after process automation")
    add_textbox(slide, "Savings reach 30% in the base case.", 1.2, 3.2, 12, 1)
    add_footer(slide, 2, "Fact table")
    out = tmp_path / "registered_fact.pptx"
    prs.save(out)
    facts = save_facts(tmp_path / "evidence.json", [{"id": "F001", "claim": "Savings", "value": 30, "unit": "%", "definition": "Base case", "source_type": "assumption", "source": "Test", "retrieved": "2026-07-05", "counter_evidence": "n/a", "used_on_pages": [2]}])
    findings = run_qa(out, facts)
    assert not any(f.check == "fact_consistency" for f in findings)


def test_unregistered_number_warns(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Fact test deck", "English cover", "Client", "2026-07-05")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Savings can reach 30% after process automation")
    add_textbox(slide, "Savings reach 30% in the base case.", 1.2, 3.2, 12, 1)
    add_footer(slide, 2, "Fact table")
    out = tmp_path / "unregistered_fact.pptx"
    prs.save(out)
    facts = save_facts(tmp_path / "evidence.json", [{"id": "F999", "claim": "Other", "value": 99, "unit": "%", "definition": "Other", "source_type": "assumption", "source": "Test", "retrieved": "2026-07-05", "counter_evidence": "n/a", "used_on_pages": []}])
    findings = run_qa(out, facts)
    assert any(f.severity == "warning" and f.check == "fact_consistency" and "Unregistered number" in f.message for f in findings)


def test_fact_value_mismatch_errors(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Fact test deck", "English cover", "Client", "2026-07-05")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Savings can reach 25% after process automation")
    add_textbox(slide, "Savings reach 25% in the base case.", 1.2, 3.2, 12, 1)
    add_footer(slide, 2, "Fact table")
    out = tmp_path / "mismatch_fact.pptx"
    prs.save(out)
    facts = save_facts(tmp_path / "evidence.json", [{"id": "F001", "claim": "Savings", "value": 30, "unit": "%", "definition": "Base case", "source_type": "assumption", "source": "Test", "retrieved": "2026-07-05", "counter_evidence": "n/a", "used_on_pages": [2]}])
    findings = run_qa(out, facts)
    assert any(f.severity == "error" and f.check == "fact_consistency" and "F001" in f.message for f in findings)
