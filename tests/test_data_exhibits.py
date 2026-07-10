from pathlib import Path
import json
import sys

from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import (
    benchmark_bar,
    cagr_annotation,
    chart_with_data_table,
    dense_table,
    driver_tree,
    native_chart,
    add_textbox,
)
from scripts.qa_pptx import run_qa


def blank_deck():
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def test_dense_table_supports_highlights_and_numeric_alignment():
    _, slide = blank_deck()
    frame = dense_table(
        slide,
        ["Metric", "FY24", "FY25", "Delta"],
        [["Revenue", 100, 120, 0.20], ["Cost", 70, 75, 0.071]],
        1, 2, 20, 5,
        numeric_columns=[1, 2, 3],
        number_formats={1: ",.0f", 2: ",.0f", 3: ".1%"},
        delta_columns=[3],
        highlight_rows=[0],
    )
    assert frame.has_table
    assert frame.table.cell(1, 0).text == "Revenue"
    assert frame.table.cell(1, 3).text == "+20.0%"


def test_driver_tree_renders_currency_prefix():
    _, slide = blank_deck()
    driver_tree(
        slide,
        {
            "label": "TCO",
            "value": 24,
            "prefix": "$",
            "unit": "m",
            "children": [
                {"label": "Platform", "value": 8, "prefix": "$", "unit": "m"},
                {"label": "Delivery", "value": 7, "prefix": "$", "unit": "m"},
            ],
        },
        x=1, y=2, w=20, h=7,
    )
    texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    assert any("$24m" in text for text in texts)
    assert any("$8m" in text for text in texts)


def test_native_chart_supports_horizontal_area_and_combo(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    native_chart(slide, "bar_h", ["A", "B"], [("Value", [10, 20])], 1, 2, 12, 6)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    native_chart(slide, "area_stacked", ["A", "B"], [("One", [10, 12]), ("Two", [4, 6])], 1, 2, 12, 6)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    frames = native_chart(
        slide,
        "combo",
        ["FY23", "FY24"],
        [("Revenue", [100, 120])],
        1, 2, 12, 6,
        secondary_series=[("Growth", [10, 20])],
        secondary_value_format='0"%"',
    )
    assert len(frames) == 2
    out = tmp_path / "charts.pptx"
    prs.save(out)
    assert out.stat().st_size > 0


def test_benchmark_and_cagr_helpers():
    _, slide = blank_deck()
    benchmark_bar(
        slide,
        ["Company", "Peer A", "Peer B"],
        [31, 42, 37],
        highlight="Company",
        benchmark_value=35,
        unit="%",
    )
    cagr = cagr_annotation(slide, 48, 105, 3, 2, 2, 10, 2)
    assert round(cagr * 100, 1) == 29.8
    texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    assert any("35" in text for text in texts)
    assert any("29.8%" in text for text in texts)


def test_chart_with_data_table_composes_both_elements():
    _, slide = blank_deck()

    def build_chart(target_slide, x, y, w, h):
        return native_chart(target_slide, "column", ["A", "B"], [("Value", [10, 20])], x, y, w, h)

    chart, table = chart_with_data_table(
        slide,
        build_chart,
        ["Metric", "A", "B"],
        [["Value", 10, 20]],
        1, 2, 20, 10,
        table_kwargs={"numeric_columns": [1, 2]},
    )
    assert chart.has_chart
    assert table.has_table


def _write_facts(path, facts):
    path.write_text(json.dumps({"facts": facts}, indent=2), encoding="utf-8")
    return path


def test_data_density_warns_below_registered_number_floor(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Density test", "Demo", "Client", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Two metrics are not enough for a research-heavy page")
    add_textbox(slide, "Savings reach 20% and payback is 18 months.", 1.2, 3.2, 20, 1)
    add_footer(slide, 2, "Test source")
    out = tmp_path / "density.pptx"
    prs.save(out)
    facts = _write_facts(
        tmp_path / "evidence.json",
        [
            {"id": "F001", "claim": "Savings", "value": 20, "unit": "%", "source_type": "assumption", "used_on_pages": [2]},
            {"id": "F002", "claim": "Payback", "value": 18, "unit": "months", "source_type": "assumption", "used_on_pages": [2]},
        ],
    )
    findings = run_qa(out, facts)
    assert any(f.check == "data_density" and f.slide == 2 for f in findings)


def test_qualitative_claim_without_number_warns(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Claim test", "Demo", "Client", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Operating redesign should improve economics")
    add_textbox(slide, "成本将显著降低。", 1.2, 3.2, 20, 1)
    add_footer(slide, 2, "Test source")
    out = tmp_path / "claim.pptx"
    prs.save(out)
    findings = run_qa(out)
    assert any(f.check == "qualitative_claim" and f.slide == 2 for f in findings)
