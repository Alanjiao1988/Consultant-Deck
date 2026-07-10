from pathlib import Path
import sys

import yaml
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import add_textbox
from scripts.qa_pptx import run_qa


def test_missing_fact_file_with_analytical_brief_is_error(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Missing facts test", "Synthetic", "Test", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Operating redesign should precede broad rollout")
    add_textbox(slide, "The page contains only conceptual analysis.", 1.2, 3.2, 20, 1)
    add_footer(slide, 2, "Synthetic test source")
    deck = tmp_path / "missing_facts.pptx"
    prs.save(deck)

    briefs = tmp_path / "briefs.yaml"
    briefs.write_text(
        yaml.safe_dump(
            {
                "content_density_target": "research-heavy",
                "pages": [
                    {"page": 1, "page_role": "cover"},
                    {"page": 2, "page_role": "core_argument"},
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    findings = run_qa(deck, briefs_path=briefs)
    assert any(
        item.severity == "error"
        and item.check == "data_density"
        and "No registered numeric facts" in item.message
        for item in findings
    )
