from pathlib import Path
import json
import sys

import pytest
import yaml
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import add_textbox
from scripts.qa_pptx import main as qa_main, run_qa


def make_deck(path: Path, title: str, body: str, *, body_box=(1.2, 3.2, 20, 1)) -> Path:
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "QA negative-path test", "Synthetic", "Test", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, title)
    add_textbox(slide, body, *body_box)
    add_footer(slide, 2, "Synthetic test source")
    prs.save(path)
    return path


def test_empty_fact_table_with_analytical_brief_is_error(tmp_path):
    deck = make_deck(
        tmp_path / "empty_facts.pptx",
        "Operating redesign should precede broad rollout",
        "The page contains only conceptual analysis.",
    )
    facts = tmp_path / "evidence.json"
    facts.write_text(json.dumps({"facts": [], "calculations": []}), encoding="utf-8")
    briefs = tmp_path / "briefs.yaml"
    briefs.write_text(
        yaml.safe_dump(
            {
                "content_density_target": "research-heavy",
                "pages": [
                    {"page": 1, "page_role": "cover"},
                    {
                        "page": 2,
                        "page_role": "core_argument",
                        "content_density_target": "research-heavy",
                    },
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    findings = run_qa(deck, facts, briefs)
    assert any(
        item.severity == "error"
        and item.check == "data_density"
        and "No registered numeric facts" in item.message
        for item in findings
    )
    assert qa_main([str(deck), "--facts", str(facts), "--briefs", str(briefs)]) == 1


@pytest.mark.parametrize(
    "claim",
    [
        "流程将显著优化。",
        "整体效果显著。",
        "方案大幅优于现状。",
        "团队可以快速迭代。",
        "结果明显更好。",
        "Performance improves dramatically.",
        "The result is considerably better.",
        "The target is vastly more efficient.",
        "The new design is far better.",
    ],
)
def test_broad_qualitative_escape_forms_warn(tmp_path, claim):
    deck = make_deck(
        tmp_path / "claim.pptx",
        "Operating redesign should improve delivery discipline",
        claim,
    )
    findings = run_qa(deck)
    assert any(item.check == "qualitative_claim" and item.slide == 2 for item in findings)


def test_unsupported_qualitative_title_is_scanned(tmp_path):
    deck = make_deck(
        tmp_path / "title_claim.pptx",
        "大幅提升运营效率",
        "The body describes the operating model.",
    )
    findings = run_qa(deck)
    assert any(item.check == "qualitative_claim" and item.slide == 2 for item in findings)


def test_text_overflow_heuristic_warns(tmp_path):
    deck = make_deck(
        tmp_path / "overflow.pptx",
        "The operating model should improve delivery discipline",
        "This deliberately long sentence is placed in a very small box so that the conservative text overflow heuristic must identify it.",
        body_box=(1.2, 3.2, 3.0, 0.35),
    )
    findings = run_qa(deck)
    assert any(item.check == "text_overflow" and item.slide == 2 for item in findings)


def test_independent_textbox_overlap_warns(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Overlap test", "Synthetic", "Test", "2026-07-10")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "The operating model should improve delivery discipline")
    add_textbox(slide, "First independent text box", 2.0, 4.0, 8.0, 1.2)
    add_textbox(slide, "Second independent text box", 3.0, 4.2, 8.0, 1.2)
    add_footer(slide, 2, "Synthetic test source")
    deck = tmp_path / "overlap.pptx"
    prs.save(deck)

    findings = run_qa(deck)
    assert any(item.check == "text_overlap" and item.slide == 2 for item in findings)


def test_success_message_lists_enabled_modules(tmp_path, capsys):
    deck = make_deck(
        tmp_path / "clean.pptx",
        "The operating model should improve delivery discipline",
        "Implementation detail follows in the next workstream.",
    )
    assert qa_main([str(deck)]) == 0
    output = capsys.readouterr().out
    assert "QA passed: 0 findings" in output
    assert "text-overflow" in output
    assert "text-overlap" in output
    assert "facts=not-supplied" in output
