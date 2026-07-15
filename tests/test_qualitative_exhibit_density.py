from __future__ import annotations

import json
from pathlib import Path
import sys

from pptx import Presentation
from pptx.util import Cm
import yaml

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import add_textbox
from scripts.qa_pptx import run_qa


def test_qualitative_exhibit_explicit_zero_uses_alternative_density_contract(tmp_path):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover_page(cover, "Regulatory map", "Qualitative exhibit test", "Client", "2026-07-15")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_action_title(slide, "Regulatory operating models differ by jurisdiction and workload")
    add_textbox(slide, "Coverage: 24 reviewed; 6 representative markets shown.", 1.2, 3.2, 20, 0.8)
    add_textbox(slide, "As of 2026-07-15. See Appendix A6 for the full reviewed universe.", 1.2, 4.2, 20, 0.8)
    add_footer(slide, 2, "Illustrative evidence table")

    pptx_path = tmp_path / "qualitative.pptx"
    prs.save(pptx_path)

    evidence_path = tmp_path / "evidence.json"
    evidence_path.write_text(
        json.dumps(
            {
                "facts": [
                    {"id": "F001", "claim": "Markets reviewed", "value": 24, "unit": "plain", "source_type": "assumption", "used_on_pages": [2]},
                    {"id": "F002", "claim": "Markets shown", "value": 6, "unit": "plain", "source_type": "assumption", "used_on_pages": [2]},
                    {"id": "F003", "claim": "US classification", "text_value": "Fragmented oversight", "source_type": "assumption", "used_on_pages": [2]},
                    {"id": "F004", "claim": "EU classification", "text_value": "Unified rulebook", "source_type": "assumption", "used_on_pages": [2]},
                ],
                "calculations": [],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    briefs_path = tmp_path / "briefs.yaml"
    briefs_path.write_text(
        yaml.safe_dump(
            {
                "content_density_target": "research-heavy",
                "pages": [
                    {
                        "page": 2,
                        "page_role": "core_argument",
                        "content_density_target": "qualitative-exhibit",
                        "min_registered_numbers": 0,
                        "exhibit_manifest": "output/deck.exhibits.json",
                        "appendix_link": "A6",
                    }
                ],
            },
            sort_keys=False,
        ),
        encoding="utf-8",
    )

    findings = run_qa(pptx_path, evidence_path, briefs_path)
    assert not any(finding.check == "data_density" and finding.slide == 2 for finding in findings)
