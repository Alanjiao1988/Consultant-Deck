"""PowerPoint-native consulting shapes for python-pptx."""
from __future__ import annotations

from math import ceil
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

PRIMARY = "1F3864"
ACCENT = "C00000"
GRAY_TEXT = "333333"
GRAY_2 = "595959"
GRAY_3 = "A6A6A6"
GRAY_LINE = "D9D9D9"
GRAY_FILL = "F2F2F2"
FONT_LATIN = "Arial"
FONT_EA = "Microsoft YaHei"


def set_font(run, size_pt=11, bold=False, color=GRAY_TEXT, italic=False):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = FONT_LATIN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_EA)


def add_textbox(slide, text, x, y, w, h, size=11, bold=False, color=GRAY_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    set_font(run, size, bold, color, italic)
    return tb


def _solid(shape, hex_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _line(shape, hex_color=None, width_pt=0.75):
    if hex_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor.from_string(hex_color)
        shape.line.width = Pt(width_pt)


def _format_number(value: Any, fmt: str | None = None) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return str(value)
    if not isinstance(value, (int, float)):
        return str(value)
    if fmt:
        try:
            return format(value, fmt)
        except (ValueError, TypeError):
            pass
    if isinstance(value, float) and not value.is_integer():
        return f"{value:,.1f}"
    return f"{value:,.0f}"


def harvey_ball(slide, x, y, fill_pct, diameter=0.55):
    d = Cm(diameter)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), d, d)
    _solid(ring, "FFFFFF")
    _line(ring, PRIMARY, 1.0)
    if fill_pct >= 100:
        _solid(ring, PRIMARY)
    return ring


def matrix_2x2(slide, x, y, size_cm, x_label, y_label, points=None, quadrant_labels=None):
    half = size_cm / 2
    for qx, qy, fill in [(0, half, "FFFFFF"), (half, half, "FFFFFF"), (0, 0, "FFFFFF"), (half, 0, GRAY_FILL)]:
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + qx), Cm(y + qy), Cm(half), Cm(half))
        _solid(r, fill)
        _line(r, GRAY_LINE)
    add_textbox(slide, x_label, x, y + size_cm + 0.15, size_cm, 0.5, 10, True, GRAY_2, PP_ALIGN.CENTER)
    add_textbox(slide, y_label, x - 0.6, y - 0.7, size_cm, 0.5, 10, True, GRAY_2)
    if quadrant_labels:
        for label, (qx, qy) in zip(quadrant_labels, [(0, half), (half, half), (0, 0), (half, 0)]):
            add_textbox(slide, label, x + qx + 0.15, y + qy + 0.1, half - 0.3, 0.5, 9, False, GRAY_3, italic=True)
    if points:
        for name, px, py in points:
            d = 0.45
            cx = x + px * size_cm - d / 2
            cy = y + (1 - py) * size_cm - d / 2
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(cx), Cm(cy), Cm(d), Cm(d))
            _solid(dot, PRIMARY)
            _line(dot, None)
            add_textbox(slide, name, cx + d + 0.08, cy - 0.05, 3.5, 0.5, 9)


def matrix_2x2_with_insights(slide, x, y, size_cm, x_label, y_label, points=None,
                             quadrant_labels=None, quadrant_implications=None, insight_title="Implications"):
    matrix_2x2(slide, x, y, size_cm, x_label, y_label, points, quadrant_labels)
    if quadrant_implications:
        ix = x + size_cm + 0.8
        add_textbox(slide, insight_title, ix, y, 10, 0.6, 11, True, PRIMARY)
        cy = y + 0.8
        for label, implication in quadrant_implications:
            add_textbox(slide, label, ix, cy, 2.8, 0.45, 9, True, PRIMARY)
            add_textbox(slide, implication, ix + 3.0, cy, 10.5, 0.8, 9.5)
            cy += 1.0


def waterfall(slide, items, x, y, w, h, unit=""):
    starts = [v for _, v, k in items if k == "start"]
    ends = [v for _, v, k in items if k == "end"]
    deltas = [v for _, v, k in items if k == "delta"]
    if starts and ends and abs(starts[0] + sum(deltas) - ends[0]) > 1e-6:
        raise ValueError("waterfall does not reconcile")
    vals = []
    run = 0
    for _, v, k in items:
        if k in ("start", "end"):
            run = v
            vals.append(v)
        else:
            vals += [run, run + v]
            run += v
    vmin, vmax = min(vals + [0]), max(vals + [0])
    span = vmax - vmin or 1
    bar_w = w / (len(items) * 1.5)
    gap = bar_w * 0.5
    zero_y = y + h - 0.5 - (0 - vmin) * ((h - 1.2) / span)
    scale = (h - 1.2) / span
    run = 0
    for i, (label, value, kind) in enumerate(items):
        bx = x + i * (bar_w + gap)
        if kind in ("start", "end"):
            top = zero_y - value * scale
            height = abs(value) * scale
            run = value
            fill = GRAY_2
        else:
            y0 = zero_y - run * scale
            run += value
            y1 = zero_y - run * scale
            top = min(y0, y1)
            height = abs(value) * scale
            fill = PRIMARY if value >= 0 else ACCENT
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(bx), Cm(top), Cm(bar_w), Cm(height))
        _solid(bar, fill)
        _line(bar, None)
        add_textbox(slide, f"{value:+g}{unit}" if kind == "delta" else f"{value:g}{unit}", bx - 0.2, top - 0.55, bar_w + 0.4, 0.5, 9, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, label, bx - 0.3, y + h - 0.4, bar_w + 0.6, 0.8, 8.5, align=PP_ALIGN.CENTER)


def native_chart(slide, chart_type, categories, series, x, y, w, h, show_legend=None, show_values=True,
                 show_axis=True, show_gridlines=False, value_format=None, highlight_series=0,
                 secondary_series=None, secondary_value_format=None):
    """Create editable native charts.

    Supported chart types: column/bar, bar_h, stacked_bar, line, area_stacked and
    combo. Combo uses two aligned native charts so the result remains editable.
    """
    if chart_type == "combo":
        if not secondary_series:
            raise ValueError("combo requires secondary_series")
        return combo_chart(
            slide, categories, series, secondary_series, x, y, w, h,
            primary_value_format=value_format,
            secondary_value_format=secondary_value_format,
            show_legend=True if show_legend is None else show_legend,
        )

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    mapping = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar_h": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    }
    if chart_type not in mapping:
        raise ValueError(f"unsupported chart_type: {chart_type}")
    if not categories or not series:
        raise ValueError("categories and series are required")
    for name, values in series:
        if len(values) != len(categories):
            raise ValueError(f"series length mismatch: {name}")

    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)

    frame = slide.shapes.add_chart(mapping[chart_type], Cm(x), Cm(y), Cm(w), Cm(h), data)
    chart = frame.chart
    palette = [PRIMARY, "6C86B3", GRAY_3, "A9B9D6"]

    for i, s in enumerate(chart.series):
        color = PRIMARY if i == highlight_series else palette[i % len(palette)]
        if chart_type == "line":
            s.format.line.color.rgb = RGBColor.from_string(color)
            s.format.line.width = Pt(2)
        else:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = RGBColor.from_string(color)
            s.format.line.fill.background()
        if show_values:
            s.data_labels.show_value = True
            s.data_labels.font.size = Pt(9)
            s.data_labels.font.name = FONT_LATIN
            if value_format:
                s.data_labels.number_format = value_format

    chart.has_legend = (len(series) > 1) if show_legend is None else bool(show_legend)
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    try:
        chart.category_axis.visible = bool(show_axis)
        chart.value_axis.visible = bool(show_axis)
        chart.value_axis.has_major_gridlines = bool(show_gridlines)
        if value_format:
            chart.value_axis.tick_labels.number_format = value_format
    except Exception:
        pass
    return frame


def combo_chart(slide, categories, column_series, line_series, x, y, w, h,
                primary_value_format=None, secondary_value_format=None, show_legend=True):
    """Create an editable column + line combination using aligned native charts.

    python-pptx does not expose a stable public API for mixed-plot charts. This
    helper overlays a transparent line chart on a column chart. It avoids
    brittle XML surgery while keeping both plots editable in PowerPoint.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.oxml.xmlchemy import OxmlElement

    if not categories or not column_series or not line_series:
        raise ValueError("categories, column_series and line_series are required")
    for name, values in list(column_series) + list(line_series):
        if len(values) != len(categories):
            raise ValueError(f"series length mismatch: {name}")

    col_data = CategoryChartData()
    col_data.categories = categories
    for name, values in column_series:
        col_data.add_series(name, values)
    col_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), col_data
    )
    col_chart = col_frame.chart
    for i, series_obj in enumerate(col_chart.series):
        series_obj.format.fill.solid()
        series_obj.format.fill.fore_color.rgb = RGBColor.from_string(
            PRIMARY if i == 0 else "6C86B3"
        )
        series_obj.format.line.fill.background()
        series_obj.data_labels.show_value = True
        series_obj.data_labels.font.size = Pt(8.5)
        if primary_value_format:
            series_obj.data_labels.number_format = primary_value_format
    col_chart.has_legend = bool(show_legend)
    if col_chart.has_legend:
        col_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        col_chart.legend.include_in_layout = False
        col_chart.legend.font.size = Pt(8.5)
    col_chart.value_axis.has_major_gridlines = True
    if primary_value_format:
        col_chart.value_axis.tick_labels.number_format = primary_value_format

    line_data = CategoryChartData()
    line_data.categories = categories
    for name, values in line_series:
        line_data.add_series(name, values)
    line_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Cm(x), Cm(y), Cm(w), Cm(h), line_data
    )
    line_chart = line_frame.chart
    for i, series_obj in enumerate(line_chart.series):
        color = ACCENT if i == 0 else GRAY_2
        series_obj.format.line.color.rgb = RGBColor.from_string(color)
        series_obj.format.line.width = Pt(2)
        series_obj.data_labels.show_value = True
        series_obj.data_labels.font.size = Pt(8.5)
        if secondary_value_format:
            series_obj.data_labels.number_format = secondary_value_format
    line_chart.has_legend = False
    line_chart.category_axis.visible = False
    line_chart.value_axis.visible = False
    line_chart.value_axis.has_major_gridlines = False

    for parent in (line_chart._chartSpace, line_chart._chartSpace.chart.plotArea):
        sp_pr = parent.find(qn("c:spPr"))
        if sp_pr is None:
            sp_pr = OxmlElement("c:spPr")
            parent.append(sp_pr)
        for child in list(sp_pr):
            if child.tag in {qn("a:solidFill"), qn("a:gradFill"), qn("a:pattFill"), qn("a:noFill")}:
                sp_pr.remove(child)
        sp_pr.insert(0, OxmlElement("a:noFill"))

    return col_frame, line_frame


def cagr_annotation(slide, start_value, end_value, periods, x1, y1, x2, y2,
                    label_prefix="CAGR", size=9.5):
    """Draw a CAGR arrow and label. Returns the calculated CAGR."""
    if start_value <= 0 or end_value <= 0 or periods <= 0:
        raise ValueError("start_value, end_value and periods must be positive")
    cagr = (end_value / start_value) ** (1 / periods) - 1
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y1), Cm(x2), Cm(y2)
    )
    connector.line.color.rgb = RGBColor.from_string(PRIMARY)
    connector.line.width = Pt(1.5)
    try:
        connector.line.end_arrowhead = True
    except Exception:
        pass
    mid_x = (x1 + x2) / 2 - 1.2
    mid_y = (y1 + y2) / 2 - 0.45
    add_textbox(slide, f"{label_prefix} {cagr:.1%}", mid_x, mid_y, 2.8, 0.45,
                size, True, PRIMARY, PP_ALIGN.CENTER)
    return cagr


def dense_table(slide, headers, rows, x, y, w, h, column_widths=None,
                number_formats=None, numeric_columns=None, highlight_rows=None,
                highlight_columns=None, total_rows=None, delta_columns=None,
                font_size=9.0):
    """Render a dense, editable consulting data table.

    `number_formats` maps column index to a Python format specifier, for example
    `{1: ",.1f", 2: ".1%"}`. `delta_columns` colour positive values in primary
    blue and negative values in accent red.
    """
    if not headers:
        raise ValueError("headers are required")
    if any(len(row) != len(headers) for row in rows):
        raise ValueError("every row must match header length")

    row_count = len(rows) + 1
    col_count = len(headers)
    frame = slide.shapes.add_table(row_count, col_count, Cm(x), Cm(y), Cm(w), Cm(h))
    table = frame.table

    if column_widths:
        if len(column_widths) != col_count:
            raise ValueError("column_widths length must match headers")
        total = sum(column_widths)
        if total <= 0:
            raise ValueError("column_widths must sum to a positive value")
        for idx, fraction in enumerate(column_widths):
            table.columns[idx].width = Cm(w * fraction / total)

    numeric_columns = set(numeric_columns or [])
    highlight_rows = set(highlight_rows or [])
    highlight_columns = set(highlight_columns or [])
    total_rows = set(total_rows or [])
    delta_columns = set(delta_columns or [])
    number_formats = number_formats or {}

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(PRIMARY)
        cell.margin_left = cell.margin_right = Cm(0.08)
        cell.margin_top = cell.margin_bottom = Cm(0.03)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if col_idx in numeric_columns else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(header)
        set_font(run, font_size, True, "FFFFFF")

    for row_idx, row in enumerate(rows, start=1):
        source_idx = row_idx - 1
        is_total = source_idx in total_rows
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            fill = "FFFFFF"
            if source_idx in highlight_rows or col_idx in highlight_columns:
                fill = "D6DEEC"
            elif row_idx % 2 == 0:
                fill = GRAY_FILL
            if is_total:
                fill = "E8EBF1"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(fill)
            cell.margin_left = cell.margin_right = Cm(0.08)
            cell.margin_top = cell.margin_bottom = Cm(0.02)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if col_idx in numeric_columns else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = _format_number(value, number_formats.get(col_idx))
            color = GRAY_TEXT
            if col_idx in delta_columns and isinstance(value, (int, float)):
                color = PRIMARY if value >= 0 else ACCENT
                if value > 0 and not run.text.startswith("+"):
                    run.text = "+" + run.text
            set_font(run, font_size, is_total or source_idx in highlight_rows, color)

    return frame


def benchmark_bar(slide, entities, values, highlight=None, x=1.4, y=3.0, w=14.0, h=9.0,
                  benchmark_value=None, benchmark_label="Industry average",
                  value_format=",.1f", unit=""):
    """Draw a horizontal benchmark bar chart using native shapes."""
    if not entities or len(entities) != len(values):
        raise ValueError("entities and values must be non-empty and equal length")
    if not all(isinstance(value, (int, float)) for value in values):
        raise ValueError("values must be numeric")

    max_value = max(max(values), benchmark_value if benchmark_value is not None else max(values), 0)
    min_value = min(min(values), 0)
    span = max_value - min_value or 1
    label_w = min(5.0, w * 0.30)
    plot_x = x + label_w
    plot_w = w - label_w - 1.7
    row_h = h / max(len(values), 1)
    zero_x = plot_x + (-min_value / span) * plot_w

    if benchmark_value is not None:
        bx = plot_x + ((benchmark_value - min_value) / span) * plot_w
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(bx), Cm(y), Cm(bx), Cm(y + h))
        line.line.color.rgb = RGBColor.from_string(ACCENT)
        line.line.width = Pt(1)
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        add_textbox(slide, f"{benchmark_label}: {_format_number(benchmark_value, value_format)}{unit}",
                    max(x, bx - 2.6), y - 0.55, 5.2, 0.45, 8.5, True, ACCENT, PP_ALIGN.CENTER)

    highlight_index = highlight if isinstance(highlight, int) else None
    if isinstance(highlight, str) and highlight in entities:
        highlight_index = entities.index(highlight)

    for idx, (entity, value) in enumerate(zip(entities, values)):
        top = y + idx * row_h + row_h * 0.18
        bar_h = row_h * 0.54
        add_textbox(slide, entity, x, top - 0.02, label_w - 0.2, bar_h, 9.0,
                    idx == highlight_index, PRIMARY if idx == highlight_index else GRAY_TEXT,
                    PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
        value_x = plot_x + ((min(value, 0) - min_value) / span) * plot_w
        value_w = abs(value) / span * plot_w
        if value_w < 0.08:
            value_w = 0.08
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(value_x), Cm(top), Cm(value_w), Cm(bar_h)
        )
        _solid(bar, PRIMARY if idx == highlight_index else GRAY_3)
        _line(bar, None)
        label_x = value_x + value_w + 0.12 if value >= 0 else value_x - 1.5
        add_textbox(slide, f"{_format_number(value, value_format)}{unit}",
                    label_x, top - 0.02, 1.5, bar_h, 8.8, True,
                    PRIMARY if idx == highlight_index else GRAY_2,
                    PP_ALIGN.LEFT if value >= 0 else PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)

    return {"zero_x": zero_x, "benchmark_value": benchmark_value}


def driver_tree(slide, root, branches=None, x=1.4, y=3.0, w=20.0, h=9.0,
                node_w=4.0, node_h=1.15):
    """Render a quantified driver tree.

    Node schema: `{"label": str, "value": number|str, "unit": str,
    "children": [node, ...]}`. `branches` is accepted for convenience and is
    attached to `root` when root has no `children`.
    """
    if not isinstance(root, dict):
        raise ValueError("root must be a mapping")
    root = dict(root)
    if branches is not None and not root.get("children"):
        root["children"] = branches

    nodes = []
    leaves = []

    def walk(node, depth, parent=None):
        item = {"node": node, "depth": depth, "parent": parent, "children": []}
        nodes.append(item)
        children = node.get("children") or []
        for child in children:
            child_item = walk(child, depth + 1, item)
            item["children"].append(child_item)
        if not children:
            leaves.append(item)
        return item

    root_item = walk(root, 0)
    max_depth = max(item["depth"] for item in nodes)
    if max_depth == 0:
        max_depth = 1
    leaf_gap = h / max(len(leaves), 1)
    for idx, leaf in enumerate(leaves):
        leaf["y"] = y + idx * leaf_gap + (leaf_gap - node_h) / 2

    def set_internal_y(item):
        if "y" in item:
            return item["y"]
        child_ys = [set_internal_y(child) for child in item["children"]]
        item["y"] = sum(child_ys) / len(child_ys)
        return item["y"]

    set_internal_y(root_item)
    depth_gap = (w - node_w) / max_depth

    for item in nodes:
        item["x"] = x + item["depth"] * depth_gap

    for item in nodes:
        if item["parent"] is None:
            continue
        parent = item["parent"]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW,
            Cm(parent["x"] + node_w), Cm(parent["y"] + node_h / 2),
            Cm(item["x"]), Cm(item["y"] + node_h / 2),
        )
        connector.line.color.rgb = RGBColor.from_string(GRAY_3)
        connector.line.width = Pt(1)

    for item in nodes:
        node = item["node"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(item["x"]), Cm(item["y"]), Cm(node_w), Cm(node_h)
        )
        is_root = item["depth"] == 0
        _solid(shape, PRIMARY if is_root else ("D6DEEC" if item["children"] else "FFFFFF"))
        _line(shape, PRIMARY if not is_root else None, 0.8)
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        label_run = p.add_run()
        label_run.text = str(node.get("label", ""))
        set_font(label_run, 8.8, True, "FFFFFF" if is_root else GRAY_TEXT)
        value = node.get("value")
        if value is not None:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            value_run = p2.add_run()
            value_text = _format_number(value, node.get("format"))
            unit = str(node.get("unit", ""))
            value_run.text = f"{value_text}{unit}"
            set_font(value_run, 9.6, True, "FFFFFF" if is_root else PRIMARY)
    return root_item


def chart_with_data_table(slide, chart_builder, table_headers, table_rows,
                          x, y, w, h, table_share=0.34, table_kwargs=None):
    """Compose a chart above a dense table using a common page region."""
    if not 0.20 <= table_share <= 0.60:
        raise ValueError("table_share must be between 0.20 and 0.60")
    gap = 0.35
    table_h = h * table_share
    chart_h = h - table_h - gap
    chart_result = chart_builder(slide, x, y, w, chart_h)
    table_result = dense_table(
        slide, table_headers, table_rows, x, y + chart_h + gap, w, table_h,
        **(table_kwargs or {})
    )
    return chart_result, table_result


def option_evaluation_table(slide, options, criteria, scores, x, y, w, h):
    nr, nc = len(options), len(criteria)
    label_w = w * 0.28
    col_w = (w - label_w) / nc
    row_h = (h - 0.8) / nr
    add_textbox(slide, "Option", x, y, label_w, 0.7, 9.5, True, GRAY_2)
    for j, c in enumerate(criteria):
        add_textbox(slide, c, x + label_w + j * col_w, y, col_w, 0.7, 9, True, GRAY_2, PP_ALIGN.CENTER)
    max_by_col = [max(row[j] for row in scores) for j in range(nc)]
    for i, opt in enumerate(options):
        ry = y + 0.8 + i * row_h
        add_textbox(slide, opt, x, ry + 0.05, label_w - 0.1, row_h, 9.3, True)
        for j, value in enumerate(scores[i]):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + label_w + j * col_w), Cm(ry), Cm(col_w - 0.05), Cm(row_h - 0.05))
            best = value == max_by_col[j]
            _solid(cell, "D6DEEC" if best else "FFFFFF")
            _line(cell, GRAY_LINE, 0.5)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(value)
            set_font(r, 9.5, best, PRIMARY if best else GRAY_TEXT)


def risk_matrix(slide, risks, x, y, w, h):
    cell_w, cell_h = w / 5, h / 5
    for i in range(5):
        for j in range(5):
            sev = i + j
            fill = "FFFFFF" if sev < 4 else (GRAY_FILL if sev < 6 else "D6DEEC")
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + i * cell_w), Cm(y + (4 - j) * cell_h), Cm(cell_w), Cm(cell_h))
            _solid(cell, fill)
            _line(cell, GRAY_LINE, 0.5)
    add_textbox(slide, "Likelihood", x, y + h + 0.15, w, 0.45, 9.5, True, GRAY_2, PP_ALIGN.CENTER)
    add_textbox(slide, "Impact", x - 0.7, y - 0.45, h, 0.45, 9.5, True, GRAY_2)
    for name, likelihood, impact, _ in risks:
        d = 0.55
        px = x + (max(1, min(5, likelihood)) - 0.5) * cell_w - d / 2
        py = y + (5 - max(1, min(5, impact)) + 0.5) * cell_h - d / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(px), Cm(py), Cm(d), Cm(d))
        _solid(dot, ACCENT if likelihood + impact >= 8 else PRIMARY)
        _line(dot, None)
        add_textbox(slide, name, px + d + 0.08, py - 0.02, 3.2, 0.45, 8.5)


def raci_matrix(slide, activities, roles, assignments, x, y, w, h):
    label_w = w * 0.34
    col_w = (w - label_w) / len(roles)
    row_h = (h - 0.8) / len(activities)
    add_textbox(slide, "Activity", x, y, label_w, 0.7, 9.5, True, GRAY_2)
    for j, role in enumerate(roles):
        add_textbox(slide, role, x + label_w + j * col_w, y, col_w, 0.7, 8.5, True, GRAY_2, PP_ALIGN.CENTER)
    for i, act in enumerate(activities):
        ry = y + 0.8 + i * row_h
        add_textbox(slide, act, x, ry + 0.05, label_w - 0.1, row_h, 9)
        for j, role in enumerate(roles):
            val = assignments.get((act, role), "")
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + label_w + j * col_w), Cm(ry), Cm(col_w - 0.05), Cm(row_h - 0.05))
            _solid(cell, "FFFFFF")
            _line(cell, GRAY_LINE, 0.5)
            if val:
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = val
                set_font(r, 10, True, PRIMARY if val in ("A", "R") else GRAY_2)
