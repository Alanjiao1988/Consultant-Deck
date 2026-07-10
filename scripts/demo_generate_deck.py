"""Generate a data-rich mixed Chinese-English consulting demo deck."""
from __future__ import annotations

from datetime import date
from pathlib import Path
import json
import sys

import yaml
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_H, SLIDE_W, add_action_title, add_footer, cover_page
from scripts.consulting_shapes import (
    add_textbox,
    benchmark_bar,
    cagr_annotation,
    dense_table,
    driver_tree,
    native_chart,
)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fact(fid, claim, value, unit, pages, source, definition=None):
    return {
        "id": fid,
        "claim": claim,
        "value": value,
        "unit": unit,
        "period": "Illustrative base case",
        "entity": "Illustrative bank",
        "definition": definition or claim,
        "source_type": "assumption",
        "source_tier": 5,
        "source": source,
        "source_date": str(date.today()),
        "retrieved": str(date.today()),
        "counter_evidence": "Actual results depend on scope, adoption, data quality and control requirements.",
        "caveat": "Illustrative demo data; not a client forecast.",
        "confidence": "illustrative",
        "used_on_pages": pages,
    }


def build_demo_evidence():
    facts = [
        _fact("F001", "Prioritized first-wave use cases", 12, "plain", [2], "Illustrative portfolio baseline"),
        _fact("F002", "Base-case payback period", 18, "months", [2, 7], "Illustrative operations model"),
        _fact("F003", "Three-year net benefit", 18, "$m", [2, 7], "Illustrative operations model"),
        _fact("F004", "Run-rate saving", 22, "%", [2], "Illustrative operations model"),
        _fact("F004A", "Investment hurdle", 24, "months", [2], "Illustrative operations model"),
        _fact("F005", "Annual workflow volume FY2022", 48, "plain", [3], "Illustrative operations baseline"),
        _fact("F006", "Annual workflow volume FY2023", 62, "plain", [3], "Illustrative operations baseline"),
        _fact("F007", "Annual workflow volume FY2024", 81, "plain", [3], "Illustrative operations baseline"),
        _fact("F008", "Annual workflow volume FY2025", 105, "plain", [3], "Illustrative operations baseline"),
        _fact("F009", "Automation rate FY2022", 8, "%", [3], "Illustrative operations baseline"),
        _fact("F010", "Automation rate FY2023", 12, "%", [3], "Illustrative operations baseline"),
        _fact("F011", "Automation rate FY2024", 18, "%", [3], "Illustrative operations baseline"),
        _fact("F012", "Automation rate FY2025", 27, "%", [3], "Illustrative operations baseline"),
        _fact("C001", "Workflow-volume CAGR FY2022-FY2025", 29.8, "%", [3], "Illustrative operations model"),
        _fact("F013", "Illustrative bank benchmark metric", 31, "%", [4], "Illustrative peer benchmark"),
        _fact("F014", "Peer A benchmark metric", 42, "%", [4], "Illustrative peer benchmark"),
        _fact("F015", "Peer B benchmark metric", 37, "%", [4], "Illustrative peer benchmark"),
        _fact("F016", "Peer C benchmark metric", 28, "%", [4], "Illustrative peer benchmark"),
        _fact("F017", "Peer median benchmark", 35, "%", [4], "Illustrative peer benchmark"),
        _fact("F018", "Three-year TCO", 24, "$m", [5], "Illustrative cost model"),
        _fact("F019", "Platform TCO", 8, "$m", [5], "Illustrative cost model"),
        _fact("F020", "Data TCO", 6, "$m", [5], "Illustrative cost model"),
        _fact("F021", "Delivery TCO", 7, "$m", [5], "Illustrative cost model"),
        _fact("F022", "Change TCO", 3, "$m", [5], "Illustrative cost model"),
        _fact("C002", "Platform plus delivery TCO", 15, "$m", [5], "Illustrative cost model"),
        _fact("C005", "Data plus change TCO", 9, "$m", [5], "Illustrative cost model"),
    ]

    use_cases = [
        ("RM copilot", 5.8, 82),
        ("Call summary", 3.2, 91),
        ("Credit memo", 4.7, 74),
        ("KYC review", 2.9, 69),
        ("Service search", 2.4, 88),
        ("Complaint triage", 1.8, 76),
        ("Customer chatbot", 3.6, 55),
        ("Marketing content", 1.2, 63),
    ]
    for idx, (name, value, readiness) in enumerate(use_cases, start=1):
        facts.append(_fact(f"F1{idx:02d}", f"{name} annual value", value, "$m", [6], "Illustrative use-case scoring"))
        facts.append(_fact(f"F2{idx:02d}", f"{name} readiness", readiness, "%", [6], "Illustrative use-case scoring"))

    sensitivity = [
        ("Downside adoption", 40, "%", "F301"),
        ("Base adoption", 60, "%", "F302"),
        ("Upside adoption", 80, "%", "F303"),
        ("Downside net benefit", 10, "$m", "F304"),
        ("Base net benefit", 18, "$m", "F305"),
        ("Upside net benefit", 26, "$m", "F306"),
        ("Downside payback", 30, "months", "F307"),
        ("Base payback", 18, "months", "F308"),
        ("Upside payback", 12, "months", "F309"),
    ]
    for claim, value, unit, fid in sensitivity:
        facts.append(_fact(fid, claim, value, unit, [7], "Illustrative sensitivity model"))

    roadmap = [
        ("Wave 1 systems", 12, "plain", "F401"),
        ("Wave 2 systems", 18, "plain", "F402"),
        ("Wave 3 systems", 24, "plain", "F403"),
        ("Wave 1 spend", 2.4, "$m", "F404"),
        ("Wave 2 spend", 3.1, "$m", "F405"),
        ("Wave 3 spend", 4.0, "$m", "F406"),
        ("Wave 1 exit gate", 70, "%", "F407"),
        ("Wave 2 exit gate", 85, "%", "F408"),
        ("Wave 3 exit gate", 95, "%", "F409"),
        ("Total systems", 54, "plain", "C003"),
        ("Total roadmap spend", 9.5, "$m", "C004"),
    ]
    for claim, value, unit, fid in roadmap:
        facts.append(_fact(fid, claim, value, unit, [8], "Illustrative delivery plan"))

    return {"facts": facts, "calculations": []}


def build_demo_briefs():
    return {
        "content_density_target": "research-heavy",
        "max_framework_share": 0.25,
        "pages": [
            {
                "page": 1,
                "page_role": "cover",
                "action_title": "Data-rich AI transformation demo",
            },
            {
                "page": 2,
                "page_role": "core_argument",
                "key_question": "What should the executive team approve?",
                "action_title": "12 个优先用例可在 18 个月回收投入，并形成 22% run-rate saving",
                "evidence_ids": ["F001", "F002", "F003", "F004", "F004A"],
                "required_data_points": ["12 prioritized use cases", "18-month payback", "$18m net benefit", "22% saving"],
                "quantification": {"baseline": "current fragmented pilots", "target": "12 use cases", "gap": "18-month payback"},
                "comparison_basis": {"type": "baseline_to_target"},
                "benchmark": {"type": "investment hurdle", "threshold": "payback below 24 months"},
                "analysis_method": "executive synthesis",
                "primary_exhibit": "four-row decision metric table",
                "insight_annotations": ["Value is concentrated in the first wave", "Payback remains inside a 24-month hurdle"],
                "decision_implication": "Approve the platform-first first wave and evidence-gated scaling",
                "data_source": ["Illustrative portfolio baseline", "Illustrative operations model"],
                "caveat": "Illustrative assumptions require client validation",
                "appendix_link": "A1-A2",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 3,
                "page_role": "core_argument",
                "key_question": "How fast is the target workload and automation base growing?",
                "action_title": "业务量从 48 增至 105，约 29.8% CAGR 支撑平台化投资",
                "evidence_ids": ["F005", "F006", "F007", "F008", "F009", "F010", "F011", "F012", "C001"],
                "required_data_points": ["four-year workflow volume", "four-year automation rate", "calculated CAGR"],
                "quantification": {"baseline": 48, "target": 105, "gap": "29.8% CAGR"},
                "comparison_basis": {"type": "historical trend", "periods": ["FY2022", "FY2023", "FY2024", "FY2025"]},
                "benchmark": {"type": "internal historical", "value": "FY2022 baseline"},
                "analysis_method": "combo trend plus CAGR",
                "primary_exhibit": "column and line combo chart",
                "insight_annotations": ["Volume more than doubles", "Automation reaches 27% by FY2025"],
                "decision_implication": "Build reusable controls before volume exceeds the current operating model",
                "data_source": ["Illustrative operations baseline", "Illustrative operations model"],
                "caveat": "Volume and automation are illustrative",
                "appendix_link": "A1",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 4,
                "page_role": "core_argument",
                "key_question": "Where does current performance sit against peers?",
                "action_title": "当前 31% 指标低于 35% peer median，差距集中在可复用流程",
                "evidence_ids": ["F013", "F014", "F015", "F016", "F017"],
                "required_data_points": ["company value", "three peer values", "peer median"],
                "quantification": {"baseline": "31%", "target": "35%+", "gap": "4 percentage points"},
                "comparison_basis": {"type": "peer benchmark"},
                "benchmark": {"entities": ["Peer A", "Peer B", "Peer C"], "value": "35% median"},
                "analysis_method": "horizontal benchmark",
                "primary_exhibit": "benchmark bar with median line",
                "insight_annotations": ["Company trails the median", "Peer dispersion indicates a reachable range"],
                "decision_implication": "Use 35% as the first operating target rather than the maximum peer value",
                "data_source": ["Illustrative peer benchmark"],
                "caveat": "Peer definitions are illustrative and not perfectly comparable",
                "appendix_link": "A1",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 5,
                "page_role": "core_argument",
                "key_question": "What drives the investment requirement?",
                "action_title": "Platform 与 delivery 合计 $15m，占三年 $24m TCO 的主要部分",
                "evidence_ids": ["F018", "F019", "F020", "F021", "F022", "C002", "C005"],
                "required_data_points": ["total TCO", "four cost branches", "platform plus delivery subtotal"],
                "quantification": {"baseline": "$24m TCO", "target": "$15m platform and delivery focus", "gap": "$9m other costs"},
                "comparison_basis": {"type": "cost decomposition"},
                "benchmark": {"type": "share of total", "value": "$15m of $24m"},
                "analysis_method": "driver tree",
                "primary_exhibit": "quantified TCO driver tree",
                "insight_annotations": ["Platform and delivery dominate spend", "Change cost remains material at $3m"],
                "decision_implication": "Govern platform scope and delivery productivity as separate value levers",
                "data_source": ["Illustrative cost model"],
                "caveat": "TCO excludes financing and tax effects",
                "appendix_link": "A1",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 6,
                "page_role": "core_argument",
                "key_question": "Which use cases should scale first?",
                "action_title": "8 个候选用例中，前 3 个兼具 $3.2m–$5.8m 价值与 74%–91% 就绪度",
                "evidence_ids": [f"F1{i:02d}" for i in range(1, 9)] + [f"F2{i:02d}" for i in range(1, 9)],
                "required_data_points": ["eight annual value estimates", "eight readiness scores", "ranked first-wave set"],
                "quantification": {"range": "$1.2m-$5.8m value and 55%-91% readiness"},
                "comparison_basis": {"type": "use-case portfolio"},
                "benchmark": {"type": "first-wave threshold", "threshold": "readiness above 70%"},
                "analysis_method": "dense portfolio table",
                "primary_exhibit": "eight-row data table",
                "insight_annotations": ["RM copilot has the highest value", "Call summary has the highest readiness"],
                "decision_implication": "Scale the top three while remediating external-channel controls",
                "data_source": ["Illustrative use-case scoring"],
                "caveat": "Scores are illustrative and require business-owner validation",
                "appendix_link": "A1",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 7,
                "page_role": "appendix",
                "key_question": "How sensitive is the business case to adoption?",
                "action_title": "采用率从 40% 提升至 80% 时，净收益由 $10m 增至 $26m",
                "evidence_ids": [f"F30{i}" for i in range(1, 10)],
                "required_data_points": ["three adoption scenarios", "three net-benefit outcomes", "three payback outcomes"],
                "quantification": {"range": "40%-80% adoption, $10m-$26m net benefit"},
                "comparison_basis": {"type": "downside_base_upside"},
                "benchmark": {"type": "base case", "value": "60% adoption"},
                "analysis_method": "scenario table",
                "primary_exhibit": "three-scenario dense table",
                "insight_annotations": ["Adoption is the largest value lever", "Downside payback extends to 30 months"],
                "decision_implication": "Tie funding releases to adoption evidence",
                "data_source": ["Illustrative sensitivity model"],
                "caveat": "Sensitivity isolates adoption and holds other assumptions constant",
                "appendix_link": "Self-contained appendix",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
            {
                "page": 8,
                "page_role": "appendix",
                "key_question": "What is the quantified delivery sequence?",
                "action_title": "三波共覆盖 54 个系统、投入 $9.5m，并把 exit gate 提升至 95%",
                "evidence_ids": [f"F40{i}" for i in range(1, 10)] + ["C003", "C004"],
                "required_data_points": ["systems per wave", "spend per wave", "exit gate per wave", "portfolio totals"],
                "quantification": {"target": "54 systems, $9.5m, 95% gate"},
                "comparison_basis": {"type": "wave progression"},
                "benchmark": {"type": "exit criteria", "threshold": "70%, 85%, 95%"},
                "analysis_method": "detailed roadmap table",
                "primary_exhibit": "wave workplan table",
                "insight_annotations": ["Scope expands after each gate", "Spend is staged rather than committed upfront"],
                "decision_implication": "Release each wave only after the prior exit gate is met",
                "data_source": ["Illustrative delivery plan"],
                "caveat": "System counts and spend are illustrative",
                "appendix_link": "Self-contained appendix",
                "content_density_target": "research-heavy",
                "unresolved_gaps": [],
            },
        ],
    }


def write_demo_evidence(path):
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(build_demo_evidence(), ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def write_demo_briefs(path):
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(yaml.safe_dump(build_demo_briefs(), allow_unicode=True, sort_keys=False), encoding="utf-8")
    return path


def generate(output=ROOT / "examples" / "demo_ai_transformation.pptx",
             evidence_output=None, briefs_output=None):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    slide = blank(prs)
    cover_page(
        slide,
        "银行 AI 转型：数据密集型咨询示范",
        "Dense tables, quantified drivers, benchmarks and editable native charts",
        "Illustrative client",
        str(date.today()),
    )

    slide = blank(prs)
    add_action_title(slide, "12 个优先用例可在 18 个月回收投入，并形成 22% run-rate saving")
    dense_table(
        slide,
        ["关键指标", "数值", "决策含义"],
        [
            ["优先用例", "12 个", "聚焦首轮可复用流程"],
            ["Payback", "18 个月", "低于 24 个月投资门槛"],
            ["3 年净收益", "$18m", "税前累计 base case"],
            ["Run-rate saving", "22%", "目标流程范围"],
        ],
        1.4, 3.0, 30.4, 8.0,
        column_widths=[0.28, 0.22, 0.50],
        highlight_rows=[1],
        font_size=10,
    )
    add_textbox(slide, "建议：批准 platform-first Wave 1，并以 adoption 与 control gate 决定后续扩展。", 1.4, 12.0, 29.5, 1.0, 11, True)
    add_footer(slide, 2, "Illustrative portfolio baseline；Illustrative operations model", lang="zh")

    slide = blank(prs)
    add_action_title(slide, "业务量从 48 增至 105，约 29.8% CAGR 支撑平台化投资")
    native_chart(
        slide,
        "combo",
        ["FY2022", "FY2023", "FY2024", "FY2025"],
        [("Workflow volume", [48, 62, 81, 105])],
        1.4, 3.0, 21.5, 10.5,
        secondary_series=[("Automation rate", [8, 12, 18, 27])],
        value_format="0",
        secondary_value_format='0"%"',
    )
    cagr_annotation(slide, 48, 105, 3, 3.0, 4.2, 18.8, 4.2)
    add_textbox(slide, "105", 24.2, 3.4, 3.0, 0.7, 15, True)
    add_textbox(slide, "FY2025 workflow volume", 24.2, 4.1, 6.2, 0.6, 9)
    add_textbox(slide, "27%", 24.2, 5.3, 3.0, 0.7, 15, True)
    add_textbox(slide, "FY2025 automation rate", 24.2, 6.0, 6.2, 0.6, 9)
    add_textbox(slide, "29.8%", 24.2, 7.2, 3.0, 0.7, 15, True)
    add_textbox(slide, "FY2022-FY2025 CAGR", 24.2, 7.9, 6.2, 0.6, 9)
    add_footer(slide, 3, "Illustrative operations baseline；team calculation", lang="zh")

    slide = blank(prs)
    add_action_title(slide, "当前 31% 指标低于 35% peer median，差距集中在可复用流程")
    benchmark_bar(
        slide,
        ["Illustrative bank", "Peer A", "Peer B", "Peer C"],
        [31, 42, 37, 28],
        highlight="Illustrative bank",
        x=1.4, y=3.2, w=22.0, h=9.0,
        benchmark_value=35,
        benchmark_label="Peer median",
        value_format=".0f",
        unit="%",
    )
    add_textbox(slide, "4 pts", 25.2, 3.4, 4.0, 0.8, 16, True)
    add_textbox(slide, "gap to peer median", 25.2, 4.2, 6.0, 0.6, 9)
    add_textbox(slide, "建议先以 35% 作为运营目标，而不是直接追赶 42% 的最高同业值。", 25.2, 5.4, 6.2, 2.0, 10)
    add_footer(slide, 4, "Illustrative peer benchmark；示意数据", lang="zh")

    slide = blank(prs)
    add_action_title(slide, "Platform 与 delivery 合计 $15m，占三年 $24m TCO 的主要部分")
    driver_tree(
        slide,
        {
            "label": "3-year TCO",
            "value": 24,
            "prefix": "$",
            "unit": "m",
            "children": [
                {"label": "Platform", "value": 8, "prefix": "$", "unit": "m"},
                {"label": "Data", "value": 6, "prefix": "$", "unit": "m"},
                {"label": "Delivery", "value": 7, "prefix": "$", "unit": "m"},
                {"label": "Change", "value": 3, "prefix": "$", "unit": "m"},
            ],
        },
        x=1.4, y=3.0, w=20.5, h=10.0, node_w=4.2, node_h=1.25,
    )
    add_textbox(slide, "$15m", 24.0, 3.5, 4.5, 0.8, 17, True)
    add_textbox(slide, "Platform + delivery", 24.0, 4.3, 6.5, 0.6, 9)
    add_textbox(slide, "$9m", 24.0, 5.6, 4.5, 0.8, 17, True)
    add_textbox(slide, "Data + change", 24.0, 6.4, 6.5, 0.6, 9)
    add_footer(slide, 5, "Illustrative cost model；团队假设", lang="zh")

    use_case_rows = [
        ["RM copilot", "$5.8m", "82%", "Internal"],
        ["Call summary", "$3.2m", "91%", "Internal"],
        ["Credit memo", "$4.7m", "74%", "Controlled"],
        ["KYC review", "$2.9m", "69%", "Controlled"],
        ["Service search", "$2.4m", "88%", "Internal"],
        ["Complaint triage", "$1.8m", "76%", "Controlled"],
        ["Customer chatbot", "$3.6m", "55%", "External"],
        ["Marketing content", "$1.2m", "63%", "External"],
    ]
    slide = blank(prs)
    add_action_title(slide, "8 个候选用例中，前 3 个兼具 $3.2m–$5.8m 价值与 74%–91% 就绪度")
    dense_table(
        slide,
        ["Use case", "Annual value", "Readiness", "Exposure"],
        use_case_rows,
        1.4, 3.0, 30.4, 11.8,
        column_widths=[0.38, 0.20, 0.18, 0.24],
        numeric_columns=[1, 2],
        highlight_rows=[0, 1, 2],
        font_size=9.2,
    )
    add_footer(slide, 6, "Illustrative use-case scoring；示意数据", lang="zh")

    slide = blank(prs)
    add_action_title(slide, "采用率从 40% 提升至 80% 时，净收益由 $10m 增至 $26m")
    dense_table(
        slide,
        ["Scenario", "Adoption", "3-year net benefit", "Payback"],
        [
            ["Downside", "40%", "$10m", "30 个月"],
            ["Base", "60%", "$18m", "18 个月"],
            ["Upside", "80%", "$26m", "12 个月"],
        ],
        1.4, 3.0, 30.4, 7.2,
        column_widths=[0.28, 0.22, 0.28, 0.22],
        highlight_rows=[1],
        numeric_columns=[1, 2, 3],
        font_size=10,
    )
    add_textbox(slide, "Adoption 是最大价值杠杆；建议把后续资金释放与 60% base-case adoption 证据绑定。", 1.4, 11.2, 30.0, 1.2, 10.5, True)
    add_footer(slide, 7, "Illustrative sensitivity model；A1", lang="zh")

    slide = blank(prs)
    add_action_title(slide, "三波共覆盖 54 个系统、投入 $9.5m，并把 exit gate 提升至 95%")
    dense_table(
        slide,
        ["Wave", "Systems", "Spend", "Exit gate", "Decision gate"],
        [
            ["Wave 1", "12 个系统", "$2.4m", "70%", "Control baseline"],
            ["Wave 2", "18 个系统", "$3.1m", "85%", "Reuse proven"],
            ["Wave 3", "24 个系统", "$4.0m", "95%", "Scale approved"],
        ],
        1.4, 3.0, 30.4, 7.5,
        column_widths=[0.16, 0.18, 0.16, 0.18, 0.32],
        highlight_rows=[0],
        numeric_columns=[1, 2, 3],
        font_size=9.8,
    )
    add_textbox(slide, "54 个系统", 2.0, 11.5, 5.0, 0.8, 16, True)
    add_textbox(slide, "$9.5m", 11.0, 11.5, 4.0, 0.8, 16, True)
    add_textbox(slide, "95% final gate", 20.0, 11.5, 5.5, 0.8, 16, True)
    add_footer(slide, 8, "Illustrative delivery plan；A2", lang="zh")

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    evidence_path = Path(evidence_output) if evidence_output else output.with_suffix(".evidence.json")
    briefs_path = Path(briefs_output) if briefs_output else output.with_suffix(".briefs.yaml")
    write_demo_evidence(evidence_path)
    write_demo_briefs(briefs_path)
    print(output)
    print(evidence_path)
    print(briefs_path)
    return output, evidence_path, briefs_path


if __name__ == "__main__":
    generate()
