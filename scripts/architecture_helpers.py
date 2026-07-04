"""Architecture and operating-model exhibits for IT/AI consulting decks."""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from .consulting_shapes import PRIMARY, GRAY_TEXT, GRAY_2, GRAY_LINE, GRAY_FILL, set_font, add_textbox, _solid, _line


def capability_map(slide, domains, x, y, w, h):
    max_caps = max(len(caps) for _, caps in domains)
    dom_w = w / len(domains)
    cap_h = (h - 0.8) / max_caps
    fills = ["FFFFFF", "F2F2F2", "D6DEEC", "A9B9D6", "6C86B3", PRIMARY]
    for i, (domain, caps) in enumerate(domains):
        dx = x + i * dom_w
        add_textbox(slide, domain, dx + 0.05, y, dom_w - 0.1, 0.6, 9.5, True, PRIMARY, PP_ALIGN.CENTER)
        for j, (cap, maturity, note) in enumerate(caps):
            cy = y + 0.8 + j * cap_h
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(dx + 0.08), Cm(cy + 0.05), Cm(dom_w - 0.16), Cm(cap_h - 0.1))
            card.adjustments[0] = 0.08
            m = max(0, min(5, int(maturity)))
            _solid(card, fills[m])
            _line(card, GRAY_LINE, 0.5)
            color = "FFFFFF" if m >= 4 else GRAY_TEXT
            add_textbox(slide, cap, dx + 0.25, cy + 0.15, dom_w - 0.5, 0.35, 8.8, True, color)
            if note:
                add_textbox(slide, note, dx + 0.25, cy + 0.52, dom_w - 0.5, 0.35, 7.5, color=color)


def layered_architecture(slide, layers, x, y, w, h):
    layer_h = h / len(layers)
    for i, (layer, components) in enumerate(layers):
        ly = y + i * layer_h
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(ly), Cm(w), Cm(layer_h - 0.08))
        _solid(band, "FFFFFF" if i % 2 == 0 else GRAY_FILL)
        _line(band, GRAY_LINE, 0.5)
        add_textbox(slide, layer, x + 0.25, ly + 0.15, 4.2, layer_h - 0.3, 9.2, True, PRIMARY)
        if components:
            gap = 0.18
            comp_w = (w - 5.1 - gap * (len(components) - 1)) / len(components)
            for j, comp in enumerate(components):
                cx = x + 4.8 + j * (comp_w + gap)
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(cx), Cm(ly + 0.18), Cm(comp_w), Cm(layer_h - 0.36))
                box.adjustments[0] = 0.08
                _solid(box, "D6DEEC" if i == 0 else "FFFFFF")
                _line(box, GRAY_LINE, 0.5)
                p = box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = comp
                set_font(r, 8.5, False, GRAY_TEXT)


def hub_spoke_operating_model(slide, hub_items, spoke_items, x, y, w, h):
    hub_w = w * 0.36
    spoke_w = (w - hub_w - 1.2) / 2
    hub_x = x + spoke_w + 0.6
    hub = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(hub_x), Cm(y), Cm(hub_w), Cm(h))
    hub.adjustments[0] = 0.08
    _solid(hub, "D6DEEC"); _line(hub, PRIMARY, 1.0)
    add_textbox(slide, "Central CoE / Hub", hub_x + 0.25, y + 0.25, hub_w - 0.5, 0.5, 11, True, PRIMARY, PP_ALIGN.CENTER)
    cy = y + 1.1
    for item in hub_items:
        add_textbox(slide, f"– {item}", hub_x + 0.4, cy, hub_w - 0.8, 0.45, 9)
        cy += 0.5
    for side, sx, title, items in [("left", x, "Business spokes", spoke_items.get("business", [])), ("right", x + spoke_w + hub_w + 1.2, "Technology spokes", spoke_items.get("technology", []))]:
        spoke = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(sx), Cm(y + 0.8), Cm(spoke_w), Cm(h - 1.6))
        spoke.adjustments[0] = 0.08
        _solid(spoke, "FFFFFF"); _line(spoke, GRAY_LINE, 0.8)
        add_textbox(slide, title, sx + 0.25, y + 1.05, spoke_w - 0.5, 0.45, 10, True, GRAY_2, PP_ALIGN.CENTER)
        cy = y + 1.65
        for item in items:
            add_textbox(slide, f"– {item}", sx + 0.35, cy, spoke_w - 0.7, 0.45, 8.7)
            cy += 0.5
        x1, x2 = (sx + spoke_w, hub_x) if side == "left" else (hub_x + hub_w, sx)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(x1), Cm(y + h / 2), Cm(x2), Cm(y + h / 2))
        line.line.color.rgb = RGBColor.from_string(GRAY_LINE)
        line.line.width = Pt(1.0)
