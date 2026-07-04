"""Reusable consulting deck layouts."""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from .consulting_shapes import PRIMARY, GRAY_TEXT, GRAY_2, GRAY_3, GRAY_FILL, set_font, add_textbox, _solid, _line

SLIDE_W = 33.87
SLIDE_H = 19.05
MARGIN_L = 1.2
MARGIN_R = 1.2
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
TITLE_Y = 1.0
TITLE_H = 1.3
KICKER_Y = 2.35
CONTENT_TOP = 3.2
CONTENT_TOP_NO_KICKER = 2.6
FOOTER_Y = 18.15


def add_action_title(slide, title, kicker=None, size=19):
    title_box = add_textbox(slide, title, MARGIN_L, TITLE_Y, CONTENT_W - 6.0, TITLE_H, size=size, bold=True, color=PRIMARY)
    kicker_box = None
    if kicker:
        kicker_box = add_textbox(slide, kicker, MARGIN_L, KICKER_Y, CONTENT_W - 6.0, 0.6, size=11, color=GRAY_2)
    return title_box, kicker_box


def content_top(has_kicker=False):
    return CONTENT_TOP if has_kicker else CONTENT_TOP_NO_KICKER


def add_footer(slide, page_num=None, source=None, lang="zh"):
    if source:
        prefix = "来源：" if lang == "zh" else "Source: "
        add_textbox(slide, f"{prefix}{source}", MARGIN_L, FOOTER_Y, 24, 0.5, size=8, color=GRAY_3)
    if page_num is not None:
        add_textbox(slide, str(page_num), SLIDE_W - 2.2, FOOTER_Y, 1.0, 0.5, size=9, align=PP_ALIGN.RIGHT, color=GRAY_3)


def cover_page(slide, title, subtitle=None, client=None, date_str=None):
    add_textbox(slide, title, MARGIN_L, 6.5, CONTENT_W, 2.5, size=30, bold=True, color=PRIMARY)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN_L, 9.2, CONTENT_W, 1.0, size=14, color=GRAY_2)
    line = slide.shapes.add_connector(1, Cm(MARGIN_L), Cm(15.8), Cm(MARGIN_L + 8), Cm(15.8))
    line.line.color.rgb = RGBColor.from_string(PRIMARY)
    line.line.width = Pt(1.5)
    meta = " | ".join(x for x in [client, date_str] if x)
    if meta:
        add_textbox(slide, meta, MARGIN_L, 16.1, CONTENT_W, 0.6, size=11, color=GRAY_2)


def section_divider(slide, section_num, section_title, sections=None, current_idx=None):
    add_textbox(slide, f"{section_num:02d}", MARGIN_L, 6.0, 6, 3, size=54, bold=True, color=PRIMARY)
    add_textbox(slide, section_title, MARGIN_L, 9.6, CONTENT_W - 4, 2, size=22, bold=True, color=GRAY_TEXT)
    if sections:
        y = 13.0
        for i, item in enumerate(sections):
            cur = current_idx is not None and i == current_idx
            add_textbox(slide, f"{i + 1:02d}  {item}", MARGIN_L, y, 14, 0.7, size=11, bold=cur, color=PRIMARY if cur else GRAY_3)
            y += 0.75


def exec_summary_block(slide, blocks, x=None, y=None, w=None):
    x = MARGIN_L if x is None else x
    y = CONTENT_TOP_NO_KICKER if y is None else y
    w = CONTENT_W if w is None else w
    label_w = 3.2
    for label, conclusion, points in blocks:
        pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(label_w), Cm(0.7))
        _solid(pill, GRAY_FILL)
        _line(pill, None)
        p = pill.text_frame.paragraphs[0]
        r = p.add_run(); r.text = label
        set_font(r, 10, True, PRIMARY)
        add_textbox(slide, conclusion, x + label_w + 0.4, y, w - label_w - 0.4, 0.8, size=12, bold=True)
        y += 0.85
        for point in points:
            add_textbox(slide, f"– {point}", x + label_w + 0.4, y, w - label_w - 0.4, 0.6, size=10.5, color=GRAY_2)
            y += 0.62
        y += 0.5
    return y
