"""Editable categorical jurisdiction map with a vector SVG base layer.

The landmass is embedded as SVG so it remains sharp at any zoom level. Jurisdiction
markers, labels, leader lines and the insight rail are native PowerPoint shapes.
The helper returns a semantic exhibit manifest that can be validated by
``scripts/qa_exhibits.py`` instead of attempting to infer business meaning from
shape colours in the finished PPTX.
"""
from __future__ import annotations

from dataclasses import dataclass
from datetime import date
import json
from pathlib import Path
from typing import Any, Iterable

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.shapes.picture import CT_Picture
from pptx.parts.image import ImagePart
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MAP_ASSET = ROOT / "assets" / "maps" / "world_equal_earth.svg"
DEFAULT_ANCHORS = ROOT / "assets" / "maps" / "jurisdiction_anchors.json"
DEFAULT_THEME = ROOT / "assets" / "theme.json"

FALLBACK_COLORS = {
    "background": "FFFFFF",
    "primary": "1F3864",
    "text": "333333",
    "secondary_text": "595959",
    "light_text": "A6A6A6",
    "line": "D9D9D9",
    "fill": "F2F2F2",
    "map_land": "E8EBF1",
    "map_border": "D9D9D9",
    "category_palette": ["1F77B4", "6A3D9A", "FF7F0E", "0F8B8D", "2E8B57", "B23A48"],
}


@dataclass(frozen=True)
class LabelBox:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    def overlap_area(self, other: "LabelBox") -> float:
        width = min(self.right, other.right) - max(self.left, other.left)
        height = min(self.bottom, other.bottom) - max(self.top, other.top)
        return max(0.0, width) * max(0.0, height)


def _hex(value: str) -> str:
    return str(value).strip().lstrip("#").upper()


def load_map_theme(path: str | Path | None = None) -> dict[str, Any]:
    """Load map tokens while retaining backwards compatibility with old themes."""
    theme_path = Path(path) if path else DEFAULT_THEME
    payload: dict[str, Any] = {}
    if theme_path.exists():
        payload = json.loads(theme_path.read_text(encoding="utf-8"))
    colors = payload.get("colors", {}) if isinstance(payload, dict) else {}
    map_tokens = payload.get("map", {}) if isinstance(payload, dict) else {}
    merged = dict(FALLBACK_COLORS)
    for key in ("primary", "text", "secondary_text", "light_text", "line", "fill"):
        if colors.get(key):
            merged[key] = _hex(colors[key])
    merged["background"] = _hex(map_tokens.get("background", merged["background"]))
    merged["map_land"] = _hex(map_tokens.get("land", merged["map_land"]))
    merged["map_border"] = _hex(map_tokens.get("border", merged["map_border"]))
    palette = map_tokens.get("category_palette")
    if isinstance(palette, list) and palette:
        merged["category_palette"] = [_hex(item) for item in palette]
    fonts = payload.get("fonts", {}) if isinstance(payload, dict) else {}
    merged["font_latin"] = fonts.get("latin", "Arial")
    merged["font_ea"] = fonts.get("east_asian", "Microsoft YaHei")
    return merged


def load_anchor_registry(path: str | Path | None = None) -> dict[str, Any]:
    registry_path = Path(path) if path else DEFAULT_ANCHORS
    payload = json.loads(registry_path.read_text(encoding="utf-8"))
    anchors = payload.get("anchors")
    if not isinstance(anchors, dict) or not anchors:
        raise ValueError("anchor registry must contain a non-empty 'anchors' mapping")
    return payload


def _set_font(run: Any, theme: dict[str, Any], size: float, bold: bool = False,
              color: str | None = None, italic: bool = False) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(_hex(color or theme["text"]))
    run.font.name = theme["font_latin"]
    r_pr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = r_pr.makeelement(qn("a:ea"), {})
        r_pr.append(ea)
    ea.set("typeface", theme["font_ea"])


def _solid(shape: Any, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(_hex(color))


def _line(shape: Any, color: str | None, width: float = 0.8) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor.from_string(_hex(color))
        shape.line.width = Pt(width)


def _add_svg_picture(slide: Any, svg_bytes: bytes, x: float, y: float, w: float, h: float,
                     filename: str = "world_equal_earth.svg") -> Any:
    """Embed SVG directly in the PPTX package without rasterising it."""
    package = slide.part.package
    partname = package.next_image_partname("svg")
    image_part = ImagePart(partname, "image/svg+xml", package, svg_bytes, filename)
    r_id = slide.part.relate_to(image_part, RT.IMAGE)
    shape_id = slide.shapes._next_shape_id
    pic = CT_Picture.new_pic(
        shape_id, "Jurisdiction map base", filename, r_id,
        Cm(x), Cm(y), Cm(w), Cm(h),
    )
    slide.shapes._spTree.insert_element_before(pic, "p:extLst")
    return pic


def _clean_categories(categories: Iterable[dict[str, Any]], theme: dict[str, Any]) -> list[dict[str, Any]]:
    cleaned: list[dict[str, Any]] = []
    ids: set[str] = set()
    palette = theme["category_palette"]
    for index, category in enumerate(categories):
        if not isinstance(category, dict):
            raise ValueError("every category must be a mapping")
        category_id = str(category.get("id", "")).strip()
        label = str(category.get("label", "")).strip()
        if not category_id or not label:
            raise ValueError("every category requires non-empty id and label")
        if category_id in ids:
            raise ValueError(f"duplicate category id: {category_id}")
        ids.add(category_id)
        item = dict(category)
        item["id"] = category_id
        item["label"] = label
        item["description"] = str(category.get("description", "")).strip()
        item["color"] = _hex(category.get("color") or palette[index % len(palette)])
        cleaned.append(item)
    if not cleaned:
        raise ValueError("at least one category is required")
    if len(cleaned) > 6:
        raise ValueError("categorical jurisdiction maps support at most six categories")
    return cleaned


def _resolve_anchor(jurisdiction: dict[str, Any], registry: dict[str, Any],
                    allow_custom_anchor: bool) -> tuple[dict[str, Any], bool]:
    jurisdiction_id = str(jurisdiction.get("id", "")).strip()
    anchors = registry["anchors"]
    if jurisdiction_id in anchors:
        return dict(anchors[jurisdiction_id]), False
    custom = jurisdiction.get("anchor")
    if not allow_custom_anchor or not isinstance(custom, dict):
        raise ValueError(
            f"unregistered jurisdiction id '{jurisdiction_id}'; add it to the anchor registry "
            "or explicitly enable a justified custom anchor"
        )
    try:
        x = float(custom["x"])
        y = float(custom["y"])
    except (KeyError, TypeError, ValueError) as exc:
        raise ValueError(f"custom anchor for {jurisdiction_id} requires numeric x and y") from exc
    if not (0 <= x <= 1 and 0 <= y <= 1):
        raise ValueError(f"custom anchor for {jurisdiction_id} must use normalized 0-1 coordinates")
    rationale = str(jurisdiction.get("custom_anchor_rationale", "")).strip()
    if not rationale:
        raise ValueError(f"custom anchor for {jurisdiction_id} requires custom_anchor_rationale")
    return {
        "x": x,
        "y": y,
        "preferred_label": str(custom.get("preferred_label", "right")),
        "entity_type": str(jurisdiction.get("entity_type", "custom_group")),
    }, True


def _candidate_box(direction: str, px: float, py: float, width: float, height: float,
                   marker_d: float, gap: float = 0.14) -> LabelBox:
    if direction == "left":
        return LabelBox(px - marker_d / 2 - gap - width, py - height / 2, width, height)
    if direction == "top":
        return LabelBox(px - width / 2, py - marker_d / 2 - gap - height, width, height)
    if direction == "bottom":
        return LabelBox(px - width / 2, py + marker_d / 2 + gap, width, height)
    return LabelBox(px + marker_d / 2 + gap, py - height / 2, width, height)


def _inside(box: LabelBox, bounds: LabelBox) -> bool:
    return (
        box.left >= bounds.left and box.top >= bounds.top
        and box.right <= bounds.right and box.bottom <= bounds.bottom
    )


def _place_label(preferred: str, px: float, py: float, width: float, height: float,
                 marker_d: float, bounds: LabelBox, occupied: list[LabelBox]) -> tuple[LabelBox, str]:
    directions = [preferred] + [item for item in ("right", "left", "top", "bottom") if item != preferred]
    best: tuple[float, LabelBox, str] | None = None
    for direction in directions:
        box = _candidate_box(direction, px, py, width, height, marker_d)
        if not _inside(box, bounds):
            continue
        overlap = sum(box.overlap_area(existing) for existing in occupied)
        if overlap == 0:
            return box, direction
        if best is None or overlap < best[0]:
            best = (overlap, box, direction)
    if best is not None:
        return best[1], best[2]
    box = _candidate_box(preferred, px, py, width, height, marker_d)
    clamped = LabelBox(
        min(max(box.left, bounds.left), bounds.right - width),
        min(max(box.top, bounds.top), bounds.bottom - height),
        width, height,
    )
    return clamped, preferred


def _add_label(slide: Any, theme: dict[str, Any], jurisdiction: dict[str, Any],
               category: dict[str, Any], box: LabelBox) -> Any:
    shape = slide.shapes.add_textbox(Cm(box.left), Cm(box.top), Cm(box.width), Cm(box.height))
    shape.name = f"MAP_LABEL|entity={jurisdiction['id']}|category={category['id']}"
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Cm(0.04)
    tf.margin_top = tf.margin_bottom = Cm(0.02)
    first = tf.paragraphs[0]
    run = first.add_run()
    run.text = jurisdiction["label"]
    _set_font(run, theme, 8.8, True, theme["text"])
    note = jurisdiction.get("note", "")
    if note:
        second = tf.add_paragraph()
        second.space_before = Pt(0)
        run2 = second.add_run()
        run2.text = note
        _set_font(run2, theme, 7.5, False, theme["secondary_text"])
    return shape


def _leader_endpoint(box: LabelBox, direction: str) -> tuple[float, float]:
    if direction == "left":
        return box.right, box.top + box.height / 2
    if direction == "top":
        return box.left + box.width / 2, box.bottom
    if direction == "bottom":
        return box.left + box.width / 2, box.top
    return box.left, box.top + box.height / 2


def _add_legend_and_insights(slide: Any, theme: dict[str, Any], categories: list[dict[str, Any]],
                             x: float, y: float, w: float, h: float, legend_title: str,
                             coverage: dict[str, Any] | None, insights: list[str],
                             decision_implication: str | None, caveat: str, as_of: str) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    panel.name = "MAP_INSIGHT_RAIL"
    _solid(panel, theme["fill"])
    _line(panel, theme["line"], 0.7)

    cursor = y + 0.35
    title = slide.shapes.add_textbox(Cm(x + 0.35), Cm(cursor), Cm(w - 0.7), Cm(0.5))
    title.text_frame.margin_left = title.text_frame.margin_right = 0
    run = title.text_frame.paragraphs[0].add_run()
    run.text = legend_title
    _set_font(run, theme, 10.5, True, theme["primary"])
    cursor += 0.7

    item_h = min(1.35, max(0.9, (h * 0.48) / max(len(categories), 1)))
    for category in categories:
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x + 0.35), Cm(cursor + 0.08), Cm(0.36), Cm(0.36))
        marker.name = f"MAP_LEGEND|category={category['id']}"
        _solid(marker, category["color"])
        _line(marker, "FFFFFF", 0.7)
        label = slide.shapes.add_textbox(Cm(x + 0.86), Cm(cursor), Cm(w - 1.18), Cm(item_h))
        label.text_frame.clear()
        label.text_frame.word_wrap = True
        label.text_frame.margin_left = label.text_frame.margin_right = 0
        p = label.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = category["label"]
        _set_font(r, theme, 8.9, True, theme["text"])
        if category["description"]:
            p2 = label.text_frame.add_paragraph()
            p2.space_before = Pt(0)
            r2 = p2.add_run()
            r2.text = category["description"]
            _set_font(r2, theme, 7.3, False, theme["secondary_text"])
        cursor += item_h

    if coverage:
        reviewed = coverage.get("reviewed")
        shown = coverage.get("shown")
        if reviewed is not None or shown is not None:
            text = f"Coverage: {reviewed or '—'} reviewed; {shown or '—'} shown"
            tb = slide.shapes.add_textbox(Cm(x + 0.35), Cm(cursor + 0.08), Cm(w - 0.7), Cm(0.42))
            tb.text_frame.margin_left = tb.text_frame.margin_right = 0
            r = tb.text_frame.paragraphs[0].add_run()
            r.text = text
            _set_font(r, theme, 7.4, True, theme["secondary_text"])
            cursor += 0.55

    if insights:
        heading = slide.shapes.add_textbox(Cm(x + 0.35), Cm(cursor), Cm(w - 0.7), Cm(0.38))
        heading.text_frame.margin_left = heading.text_frame.margin_right = 0
        r = heading.text_frame.paragraphs[0].add_run()
        r.text = "Key implications"
        _set_font(r, theme, 8.5, True, theme["primary"])
        cursor += 0.45
        for insight in insights[:3]:
            tb = slide.shapes.add_textbox(Cm(x + 0.35), Cm(cursor), Cm(w - 0.7), Cm(0.66))
            tb.text_frame.margin_left = Cm(0.08)
            tb.text_frame.margin_right = 0
            p = tb.text_frame.paragraphs[0]
            p.text = "• " + insight
            for r in p.runs:
                _set_font(r, theme, 7.4, False, theme["text"])
            cursor += 0.66

    if decision_implication and cursor < y + h - 1.35:
        tb = slide.shapes.add_textbox(Cm(x + 0.35), Cm(cursor), Cm(w - 0.7), Cm(0.75))
        tb.text_frame.margin_left = tb.text_frame.margin_right = 0
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = "Decision: "
        _set_font(r, theme, 7.5, True, theme["primary"])
        r2 = p.add_run()
        r2.text = decision_implication
        _set_font(r2, theme, 7.5, False, theme["text"])

    footer = slide.shapes.add_textbox(Cm(x + 0.35), Cm(y + h - 0.85), Cm(w - 0.7), Cm(0.62))
    footer.text_frame.margin_left = footer.text_frame.margin_right = 0
    footer.text_frame.word_wrap = True
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"As of {as_of}. {caveat}"
    _set_font(r, theme, 6.7, False, theme["secondary_text"], italic=True)


def categorical_jurisdiction_map(
    slide: Any,
    jurisdictions: Iterable[dict[str, Any]],
    categories: Iterable[dict[str, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    classification_dimension: str,
    as_of: str,
    caveat: str,
    coverage: dict[str, Any] | None = None,
    insight_annotations: Iterable[str] | None = None,
    decision_implication: str | None = None,
    legend_title: str = "Regulatory model",
    mutually_exclusive: bool = True,
    theme_path: str | Path | None = None,
    map_asset: str | Path | None = None,
    anchor_registry: str | Path | None = None,
    render_mode: str = "svg_base_native_overlay",
    label_strategy: str = "auto",
    max_markers: int = 12,
    allow_custom_anchor: bool = False,
    page: int | None = None,
) -> dict[str, Any]:
    """Render a data-driven categorical map and return its semantic manifest."""
    dimension = str(classification_dimension).strip()
    caveat_text = str(caveat).strip()
    as_of_text = str(as_of).strip()
    if not dimension:
        raise ValueError("classification_dimension is required")
    if not caveat_text:
        raise ValueError("caveat is required for a categorical jurisdiction map")
    try:
        date.fromisoformat(as_of_text)
    except ValueError as exc:
        raise ValueError("as_of must use ISO YYYY-MM-DD format") from exc
    if render_mode != "svg_base_native_overlay":
        raise ValueError("phase 1 supports render_mode='svg_base_native_overlay' only")
    if label_strategy != "auto":
        raise ValueError("phase 1 supports label_strategy='auto' only")

    theme = load_map_theme(theme_path)
    registry = load_anchor_registry(anchor_registry)
    category_list = _clean_categories(categories, theme)
    category_by_id = {item["id"]: item for item in category_list}

    jurisdiction_list: list[dict[str, Any]] = []
    used_ids: set[str] = set()
    category_usage = {item["id"]: 0 for item in category_list}
    custom_anchor_ids: list[str] = []
    for raw in jurisdictions:
        if not isinstance(raw, dict):
            raise ValueError("every jurisdiction must be a mapping")
        item = dict(raw)
        item["id"] = str(raw.get("id", "")).strip()
        item["label"] = str(raw.get("label", "")).strip()
        item["category_id"] = str(raw.get("category_id", "")).strip()
        item["note"] = str(raw.get("note", "")).strip()
        item["entity_type"] = str(raw.get("entity_type", "country_or_market")).strip()
        evidence_ids = raw.get("evidence_ids") or []
        item["evidence_ids"] = [str(value).strip() for value in evidence_ids if str(value).strip()]
        if not item["id"] or not item["label"]:
            raise ValueError("every jurisdiction requires id and label")
        if item["id"] in used_ids:
            raise ValueError(f"duplicate jurisdiction id: {item['id']}")
        used_ids.add(item["id"])
        if item["category_id"] not in category_by_id:
            raise ValueError(f"unknown category '{item['category_id']}' for jurisdiction {item['id']}")
        if not item["note"]:
            raise ValueError(f"jurisdiction {item['id']} requires a non-empty note")
        if not item["evidence_ids"]:
            raise ValueError(f"jurisdiction {item['id']} requires at least one evidence_id")
        if item["entity_type"] in {"custom_group", "regulatory_bloc"}:
            members = [str(value).strip() for value in (raw.get("members") or []) if str(value).strip()]
            if item["entity_type"] == "custom_group" and not members:
                raise ValueError(f"custom group {item['id']} requires members")
            item["members"] = members
        anchor, custom = _resolve_anchor(raw, registry, allow_custom_anchor)
        item["_anchor"] = anchor
        if custom:
            custom_anchor_ids.append(item["id"])
        category_usage[item["category_id"]] += 1
        jurisdiction_list.append(item)

    if not jurisdiction_list:
        raise ValueError("at least one jurisdiction is required")
    if len(jurisdiction_list) > max_markers:
        raise ValueError(f"map supports at most {max_markers} jurisdiction markers")
    unused = [category_id for category_id, count in category_usage.items() if count == 0]
    if unused:
        raise ValueError(f"legend categories are not referenced: {', '.join(unused)}")
    if coverage:
        reviewed = coverage.get("reviewed")
        shown = coverage.get("shown", len(jurisdiction_list))
        if reviewed is not None and int(reviewed) < int(shown):
            raise ValueError("coverage.reviewed cannot be lower than coverage.shown")

    rail_w = max(6.2, w * 0.29)
    gap = 0.45
    map_region = LabelBox(x, y, w - rail_w - gap, h)
    rail_x = x + map_region.width + gap

    asset_path = Path(map_asset) if map_asset else DEFAULT_MAP_ASSET
    svg_text = asset_path.read_text(encoding="utf-8")
    svg_text = svg_text.replace("#E8EBF1", f"#{theme['map_land']}")
    svg_text = svg_text.replace("#D9D9D9", f"#{theme['map_border']}")
    aspect = 1000 / 520
    draw_w = map_region.width
    draw_h = draw_w / aspect
    if draw_h > map_region.height:
        draw_h = map_region.height
        draw_w = draw_h * aspect
    draw_x = map_region.left + (map_region.width - draw_w) / 2
    draw_y = map_region.top + (map_region.height - draw_h) / 2
    _add_svg_picture(slide, svg_text.encode("utf-8"), draw_x, draw_y, draw_w, draw_h, asset_path.name)

    marker_d = 0.44
    occupied: list[LabelBox] = []
    label_bounds = LabelBox(map_region.left, map_region.top, map_region.width, map_region.height)
    for item in jurisdiction_list:
        anchor = item["_anchor"]
        px = draw_x + float(anchor["x"]) * draw_w
        py = draw_y + float(anchor["y"]) * draw_h
        category = category_by_id[item["category_id"]]
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Cm(px - marker_d / 2), Cm(py - marker_d / 2),
            Cm(marker_d), Cm(marker_d),
        )
        marker.name = f"MAP_MARKER|entity={item['id']}|category={category['id']}|evidence={','.join(item['evidence_ids'])}"
        _solid(marker, category["color"])
        _line(marker, "FFFFFF", 0.9)

        label_w = 3.15 if len(item["note"]) <= 25 else 3.65
        label_h = 0.82 if item["note"] else 0.44
        preferred = str(anchor.get("preferred_label", "right"))
        box, direction = _place_label(preferred, px, py, label_w, label_h, marker_d, label_bounds, occupied)
        occupied.append(box)
        endpoint_x, endpoint_y = _leader_endpoint(box, direction)
        leader = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(px), Cm(py), Cm(endpoint_x), Cm(endpoint_y))
        leader.name = f"MAP_LEADER|entity={item['id']}"
        leader.line.color.rgb = RGBColor.from_string(theme["light_text"])
        leader.line.width = Pt(0.65)
        _add_label(slide, theme, item, category, box)

    _add_legend_and_insights(
        slide, theme, category_list, rail_x, y, rail_w, h, legend_title,
        coverage,
        [str(value).strip() for value in (insight_annotations or []) if str(value).strip()],
        str(decision_implication).strip() if decision_implication else None,
        caveat_text, as_of_text,
    )

    manifest_jurisdictions = []
    for item in jurisdiction_list:
        public_item = {key: value for key, value in item.items() if not key.startswith("_")}
        manifest_jurisdictions.append(public_item)
    return {
        "page": page,
        "type": "categorical_jurisdiction_map",
        "classification_dimension": dimension,
        "mutually_exclusive": bool(mutually_exclusive),
        "as_of": as_of_text,
        "render_mode": render_mode,
        "categories": category_list,
        "jurisdictions": manifest_jurisdictions,
        "coverage": coverage or {"reviewed": len(jurisdiction_list), "shown": len(jurisdiction_list)},
        "caveat": caveat_text,
        "insight_annotations": [str(value).strip() for value in (insight_annotations or []) if str(value).strip()],
        "decision_implication": str(decision_implication).strip() if decision_implication else None,
        "custom_anchor_ids": custom_anchor_ids,
    }


def write_exhibit_manifest(path: str | Path, exhibits: Iterable[dict[str, Any]]) -> Path:
    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(
        json.dumps({"exhibits": list(exhibits)}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return output
