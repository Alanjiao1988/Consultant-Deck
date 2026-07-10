"""Automated QA for consulting-style PPTX files."""
from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation

SLIDE_W_CM = 33.87
SLIDE_H_CM = 19.05
EMU_PER_CM = 360000

EN_ACTION_SIGNALS = {
    "can", "will", "should", "must", "need", "needs", "require", "requires",
    "enable", "enables", "offer", "offers", "provide", "provides", "prioritize",
    "prioritizes", "reduce", "reduces", "increase", "increases", "improve",
    "improves", "deliver", "delivers", "prevent", "prevents", "prefer",
    "maximize", "maximizes", "unlock", "unlocks", "shorten", "shortens",
    "reach", "reaches", "drive", "drives", "create", "creates", "balance",
    "outperform", "mitigate", "because", "within", "while", "before", "after",
    "therefore", "best", "better", "highest", "lowest", "faster", "slower",
    "more", "less", "first-wave", "platform-first",
}
CN_ACTION_SIGNALS = [
    "降低", "提升", "实现", "建议", "需要", "应", "将", "可", "必须", "能够",
    "优先", "集中", "更", "最佳", "风险", "收益", "回收", "减少", "增加",
]
BANNED_TERMS = {
    "organization": "组织",
    "strategy": "战略",
    "stakeholder": "相关方",
    "alignment": "对齐",
    "capability": "能力",
    "initiative": "举措",
}
EXEMPT_ROLES = {"cover", "section_divider", "divider", "navigation", "agenda"}
CONCEPT_ROLES = {"conceptual_framework", "concept", "framework"}
QUALITATIVE_PATTERNS = {
    # Warning-level recall is intentionally broad: the same statement must contain
    # a number, range or threshold when these strength signals are used.
    "显著": re.compile(r"显著"),
    "大幅": re.compile(r"大幅"),
    "明显": re.compile(r"明显"),
    "快速": re.compile(r"快速"),
    "大量": re.compile(r"大量"),
    "领先/优于": re.compile(r"(?:领先|优于|远高于|远低于|更好|更高|更低)"),
    "significantly": re.compile(r"\bsignificantly\b", re.IGNORECASE),
    "substantially": re.compile(r"\bsubstantially\b", re.IGNORECASE),
    "rapidly": re.compile(r"\brapidly\b", re.IGNORECASE),
    "materially": re.compile(r"\bmaterially\b", re.IGNORECASE),
    "dramatically": re.compile(r"\bdramatically\b", re.IGNORECASE),
    "considerably": re.compile(r"\bconsiderably\b", re.IGNORECASE),
    "vastly": re.compile(r"\bvastly\b", re.IGNORECASE),
    "far/much comparative": re.compile(
        r"\b(?:far|much)\s+(?:better|worse|higher|lower|faster|slower)\b",
        re.IGNORECASE,
    ),
    "leading": re.compile(r"\bleading\s+(?:peers?|the industry|the market)\b", re.IGNORECASE),
}


@dataclass
class Finding:
    severity: str
    slide: int | None
    check: str
    message: str


@dataclass
class NumberToken:
    slide: int
    value: float
    unit: str
    text: str


def cm(value: int) -> float:
    return value / EMU_PER_CM


def shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    if getattr(shape, "has_table", False):
        return "\n".join(
            cell.text
            for row in shape.table.rows
            for cell in row.cells
            if cell.text
        ).strip()
    return ""


def has_cjk(text: str) -> bool:
    return re.search(r"[\u3400-\u9fff]", text) is not None


def first_top_shape(slide: Any) -> Any | None:
    candidates = []
    for shape in slide.shapes:
        text = shape_text(shape)
        if text and cm(shape.top) <= 2.8 and cm(shape.left) <= 22:
            candidates.append((cm(shape.top), cm(shape.left), shape))
    if not candidates:
        return None
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def first_top_text(slide: Any) -> str | None:
    shape = first_top_shape(slide)
    return shape_text(shape) if shape is not None else None


def is_cover_or_divider(index: int, texts: list[str]) -> bool:
    return index == 1 or any(re.fullmatch(r"\d{2}", text.strip()) for text in texts)


def looks_like_topic_label(title: str) -> bool:
    cleaned = re.sub(r"\s+", " ", title.strip())
    if not cleaned:
        return True
    if re.search(r"\d", cleaned):
        return False
    if has_cjk(cleaned):
        if any(signal in cleaned for signal in CN_ACTION_SIGNALS):
            return False
        cjk_chars = len(re.findall(r"[\u3400-\u9fff]", cleaned))
        return cjk_chars <= 14 and len(cleaned) <= 28
    words = re.findall(r"[A-Za-z][A-Za-z\-]*", cleaned.lower())
    if any(word in EN_ACTION_SIGNALS for word in words):
        return False
    if re.search(r"[:;—–]|\b(to|by|because|through|within|without)\b", cleaned, re.IGNORECASE):
        return False
    return len(words) <= 7


def has_font_xml(path: Path, font: str) -> bool:
    marker = 'typeface="' + font + '"'
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                if marker in archive.read(name).decode("utf-8", errors="ignore"):
                    return True
    return False


def canonical_unit(unit: str | None) -> str:
    normalized = (unit or "").strip().lower()
    aliases = {
        "%": {"%", "percent", "percentage"},
        "$m": {"$m", "usd m", "usd million", "million usd"},
        "$bn": {"$bn", "usd bn", "usd billion", "billion usd"},
        "$k": {"$k", "usd k", "usd thousand"},
        "$": {"$", "usd"},
        "months": {"个月", "月", "mo", "mos", "month", "months"},
        "years": {"年", "yr", "yrs", "year", "years"},
    }
    for canonical, values in aliases.items():
        if normalized in values:
            return canonical
    return normalized or "plain"


def load_facts(path: Path | None) -> list[dict[str, Any]]:
    if path is None:
        return []
    data = json.loads(path.read_text(encoding="utf-8"))
    records = []
    for section in ("facts", "calculations"):
        for item in data.get(section, []):
            try:
                value = float(item.get("value"))
            except (TypeError, ValueError):
                continue
            records.append({
                "id": str(item.get("id", "UNKNOWN")),
                "claim": str(item.get("claim", "")),
                "value": value,
                "unit": canonical_unit(item.get("unit", "plain")),
                "source_type": str(item.get("source_type", "")),
                "used_on_pages": {int(page) for page in item.get("used_on_pages", [])},
            })
    return records


def load_briefs(path: Path | None) -> tuple[dict[str, Any], dict[int, dict[str, Any]]]:
    if path is None:
        return {}, {}
    data = yaml.safe_load(path.read_text(encoding="utf-8"))
    if isinstance(data, list):
        meta, pages = {}, data
    elif isinstance(data, dict):
        meta, pages = data, data.get("pages", [])
    else:
        return {}, {}
    mapping = {}
    for index, page in enumerate(pages, start=1):
        if not isinstance(page, dict):
            continue
        raw_page = page.get("page", page.get("page_num", index))
        try:
            mapping[int(raw_page)] = page
        except (TypeError, ValueError):
            continue
    return meta, mapping


def values_equal(left: float, right: float) -> bool:
    return abs(float(left) - float(right)) < 1e-6


def unit_compatible(left: str, right: str) -> bool:
    return canonical_unit(left) == canonical_unit(right)


def strip_dates(text: str) -> str:
    text = re.sub(r"\b20\d{2}[-/]\d{1,2}[-/]\d{1,2}\b", " ", text)
    return re.sub(r"20\d{2}年\d{1,2}月\d{1,2}日", " ", text)


def extract_number_tokens(text: str, slide_num: int) -> list[NumberToken]:
    working = strip_dates(text)
    tokens: list[NumberToken] = []

    def consume(pattern: str, unit_getter, minimum: float | None = None) -> None:
        nonlocal working
        matches = list(re.finditer(pattern, working, flags=re.IGNORECASE))
        chars = list(working)
        for match in matches:
            value = float(match.group("value").replace(",", ""))
            unit = canonical_unit(unit_getter(match))
            if minimum is None or abs(value) >= minimum:
                tokens.append(NumberToken(slide_num, value, unit, match.group(0)))
            for offset in range(match.start(), match.end()):
                chars[offset] = " "
        working = "".join(chars)

    consume(r"\$\s*(?P<value>-?\d[\d,]*(?:\.\d+)?)\s*(?P<suffix>m|bn|k)?\b", lambda m: "$" + (m.group("suffix") or ""))
    consume(r"(?P<value>-?\d[\d,]*(?:\.\d+)?)\s*%", lambda _m: "%")
    consume(r"(?P<value>-?\d[\d,]*(?:\.\d+)?)\s*(?P<unit>个月|月|months?|mo|mos)(?![A-Za-z])", lambda m: m.group("unit"))
    consume(r"(?P<value>-?\d[\d,]*(?:\.\d+)?)\s*(?P<unit>年|years?|yrs?)(?![A-Za-z])", lambda m: m.group("unit"), 6)

    for match in re.finditer(r"\b(?P<value>-?\d[\d,]*(?:\.\d+)?)\b", working):
        value = float(match.group("value").replace(",", ""))
        if abs(value) > 12:
            tokens.append(NumberToken(slide_num, value, "plain", match.group(0)))
    return tokens


def text_shapes_for_fact_scan(slide: Any, exclude_shape: Any | None = None):
    for shape in slide.shapes:
        if exclude_shape is not None and shape == exclude_shape:
            continue
        text = shape_text(shape)
        if not text:
            continue
        left, top = cm(shape.left), cm(shape.top)
        if left > 30 and top > 17:
            continue
        if top > 16.2:
            continue
        yield text


def chart_values(slide: Any) -> list[float]:
    values = []
    for shape in slide.shapes:
        if not getattr(shape, "has_chart", False):
            continue
        for series in shape.chart.series:
            for value in series.values:
                if value is not None:
                    values.append(float(value))
    return values


def token_matches_fact(token: NumberToken, fact: dict[str, Any]) -> bool:
    return unit_compatible(token.unit, fact["unit"]) and values_equal(token.value, fact["value"])


def fact_applies_to_page(fact: dict[str, Any], page: int) -> bool:
    return not fact["used_on_pages"] or page in fact["used_on_pages"]


def run_fact_consistency(prs: Presentation, facts: list[dict[str, Any]]) -> list[Finding]:
    if not facts:
        return []
    findings: list[Finding] = []
    tokens: list[NumberToken] = []
    per_page_charts: dict[int, list[float]] = {}
    per_page_text: dict[int, str] = {}

    for index, slide in enumerate(prs.slides, start=1):
        texts = [shape_text(shape) for shape in slide.shapes if shape_text(shape)]
        if is_cover_or_divider(index, texts):
            continue
        body_parts = list(text_shapes_for_fact_scan(slide))
        per_page_text[index] = "\n".join(body_parts)
        per_page_charts[index] = chart_values(slide)
        for text in body_parts:
            tokens.extend(extract_number_tokens(text, index))

    for token in tokens:
        matches = [fact for fact in facts if token_matches_fact(token, fact)]
        if not matches:
            findings.append(Finding("warning", token.slide, "fact_consistency", f"Unregistered number: {token.text}"))
        elif not any(fact_applies_to_page(fact, token.slide) for fact in matches):
            findings.append(Finding("warning", token.slide, "fact_consistency", f"Registered number appears on unlisted page: {token.text}"))

    for fact in facts:
        for page in fact["used_on_pages"]:
            page_tokens = [token for token in tokens if token.slide == page and unit_compatible(token.unit, fact["unit"])]
            if any(values_equal(token.value, fact["value"]) for token in page_tokens):
                continue
            if any(values_equal(value, fact["value"]) for value in per_page_charts.get(page, [])):
                continue
            expected_text = f"{fact['value']:g}"
            raw_pattern = re.compile(rf"(?<![\d.]){re.escape(expected_text)}(?![\d.])")
            if raw_pattern.search(per_page_text.get(page, "")):
                continue
            if page_tokens:
                seen = ", ".join(token.text for token in page_tokens[:5])
                findings.append(Finding(
                    "error", page, "fact_consistency",
                    f"Fact {fact['id']} expected {fact['value']:g}{fact['unit']} but page shows {seen}",
                ))
    return findings


def registered_fact_ids_on_slide(slide: Any, slide_num: int, facts: list[dict[str, Any]]) -> set[str]:
    title_shape = first_top_shape(slide)
    tokens: list[NumberToken] = []
    for text in text_shapes_for_fact_scan(slide, exclude_shape=title_shape):
        tokens.extend(extract_number_tokens(text, slide_num))
    chart_numbers = chart_values(slide)
    matched = set()
    for fact in facts:
        if not fact_applies_to_page(fact, slide_num):
            continue
        if any(token_matches_fact(token, fact) for token in tokens):
            matched.add(fact["id"])
        elif any(values_equal(value, fact["value"]) for value in chart_numbers):
            matched.add(fact["id"])
    return matched


def brief_role(brief: dict[str, Any]) -> str:
    return str(brief.get("page_role", brief.get("role", "core_argument"))).strip().lower()


def brief_density(brief: dict[str, Any], meta: dict[str, Any]) -> str:
    value = brief.get("content_density_target", meta.get("content_density_target", "research-heavy"))
    return str(value).strip().lower().replace("_", "-")


def density_threshold(brief: dict[str, Any], meta: dict[str, Any]) -> int:
    explicit = brief.get("min_registered_numbers")
    if explicit is not None:
        try:
            return max(0, int(explicit))
        except (TypeError, ValueError):
            pass
    return {
        "executive": 1,
        "executive-brief": 1,
        "standard": 2,
        "research-heavy": 3,
    }.get(brief_density(brief, meta), 3)


def is_density_exempt(brief: dict[str, Any]) -> bool:
    role = brief_role(brief)
    return role in EXEMPT_ROLES or (role in CONCEPT_ROLES and bool(brief.get("explicitly_requested")))


def run_data_density(prs: Presentation, facts: list[dict[str, Any]], meta: dict[str, Any], briefs: dict[int, dict[str, Any]]) -> list[Finding]:
    analytical_pages = [page for page, brief in briefs.items() if not is_density_exempt(brief)]
    if not facts:
        if analytical_pages:
            pages = ", ".join(str(page) for page in sorted(analytical_pages))
            return [Finding(
                "error", None, "data_density",
                "No registered numeric facts are available for non-exempt analytical "
                f"pages ({pages}); an empty or missing fact table cannot bypass density QA",
            )]
        return []
    findings: list[Finding] = []
    conceptual_pages = 0
    eligible_pages = 0

    for index, slide in enumerate(prs.slides, start=1):
        texts = [shape_text(shape) for shape in slide.shapes if shape_text(shape)]
        if is_cover_or_divider(index, texts):
            continue
        brief = briefs.get(index, {})
        role = brief_role(brief)
        if role in CONCEPT_ROLES and bool(brief.get("explicitly_requested")):
            conceptual_pages += 1
        if role not in EXEMPT_ROLES:
            eligible_pages += 1
        if is_density_exempt(brief):
            continue
        threshold = density_threshold(brief, meta)
        matched = registered_fact_ids_on_slide(slide, index, facts)
        if len(matched) < threshold:
            findings.append(Finding(
                "warning", index, "data_density",
                f"Only {len(matched)} registered numeric facts are visible in the slide body; expected at least {threshold}",
            ))

    if briefs and eligible_pages:
        try:
            max_share = float(meta.get("max_framework_share", 0.25))
        except (TypeError, ValueError):
            max_share = 0.25
        share = conceptual_pages / eligible_pages
        if share > max_share:
            findings.append(Finding(
                "warning", None, "data_density",
                f"Explicit conceptual-framework pages are {share:.0%} of eligible pages; maximum is {max_share:.0%}",
            ))
    return findings


def run_qualitative_claim_scan(slide: Any, slide_num: int) -> list[Finding]:
    findings: list[Finding] = []
    # Scan titles as well as body text. Brief-level title quantification is valuable,
    # but final-file QA must still catch a strong unsupported title on its own.
    body = "\n".join(text_shapes_for_fact_scan(slide))
    numeric_signal = re.compile(r"\d|%|\$|€|£|倍|点|bps?\b", re.IGNORECASE)
    seen = set()
    for sentence in re.split(r"[\n。！？!?;；]+", body):
        sentence = sentence.strip()
        if not sentence or numeric_signal.search(sentence):
            continue
        for label, pattern in QUALITATIVE_PATTERNS.items():
            if pattern.search(sentence) and label not in seen:
                findings.append(Finding(
                    "warning", slide_num, "qualitative_claim",
                    f"Qualitative claim '{label}' has no number or range in the same statement",
                ))
                seen.add(label)
    return findings


def _font_size_pt(paragraph: Any, default: float = 11.0) -> float:
    sizes = []
    for run in paragraph.runs:
        size = getattr(run.font, "size", None)
        if size is not None:
            sizes.append(float(size.pt))
    return max(sizes) if sizes else default


def _weighted_text_width(text: str, font_size_pt: float) -> float:
    units = 0.0
    for char in text:
        if char.isspace():
            units += 0.30
        elif re.match(r"[\u3400-\u9fff]", char):
            units += 1.0
        elif char.isalnum():
            units += 0.56
        else:
            units += 0.45
    return units * font_size_pt


def _text_frame_may_overflow(text_frame: Any, width_emu: int, height_emu: int) -> bool:
    if not text_frame or not any(paragraph.text.strip() for paragraph in text_frame.paragraphs):
        return False
    auto_size = getattr(text_frame, "auto_size", None)
    if auto_size is not None and "TEXT_TO_FIT_SHAPE" in str(auto_size):
        return False

    margin_left = int(getattr(text_frame, "margin_left", 0) or 0)
    margin_right = int(getattr(text_frame, "margin_right", 0) or 0)
    margin_top = int(getattr(text_frame, "margin_top", 0) or 0)
    margin_bottom = int(getattr(text_frame, "margin_bottom", 0) or 0)
    usable_width_pt = max(1.0, (width_emu - margin_left - margin_right) / 12700)
    usable_height_pt = max(1.0, (height_emu - margin_top - margin_bottom) / 12700)

    estimated_height = 0.0
    for paragraph in text_frame.paragraphs:
        text = paragraph.text or " "
        font_size = _font_size_pt(paragraph)
        wrapped_lines = max(1, math.ceil(_weighted_text_width(text, font_size) / usable_width_pt))
        estimated_height += wrapped_lines * font_size * 1.22
    return estimated_height > usable_height_pt * 1.15


def run_text_overflow_scan(slide: Any, slide_num: int) -> list[Finding]:
    findings: list[Finding] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            if _text_frame_may_overflow(shape.text_frame, shape.width, shape.height):
                findings.append(Finding(
                    "warning", slide_num, "text_overflow",
                    f"Text may overflow shape '{getattr(shape, 'name', 'unnamed')}' (heuristic estimate)",
                ))
        if getattr(shape, "has_table", False):
            seen_cells = set()
            for row_index, row in enumerate(shape.table.rows):
                for col_index, cell in enumerate(row.cells):
                    cell_key = id(cell._tc)
                    if cell_key in seen_cells:
                        continue
                    seen_cells.add(cell_key)
                    width = shape.table.columns[col_index].width
                    height = row.height
                    if _text_frame_may_overflow(cell.text_frame, width, height):
                        findings.append(Finding(
                            "warning", slide_num, "text_overflow",
                            f"Table cell R{row_index + 1}C{col_index + 1} may overflow (heuristic estimate)",
                        ))
    return findings


def _intersection_ratio(first: Any, second: Any) -> float:
    first_left, first_top = cm(first.left), cm(first.top)
    second_left, second_top = cm(second.left), cm(second.top)
    first_right, first_bottom = first_left + cm(first.width), first_top + cm(first.height)
    second_right, second_bottom = second_left + cm(second.width), second_top + cm(second.height)
    overlap_width = min(first_right, second_right) - max(first_left, second_left)
    overlap_height = min(first_bottom, second_bottom) - max(first_top, second_top)
    if overlap_width <= 0.10 or overlap_height <= 0.10:
        return 0.0
    intersection = overlap_width * overlap_height
    smaller_area = min(cm(first.width) * cm(first.height), cm(second.width) * cm(second.height))
    return intersection / smaller_area if smaller_area > 0 else 0.0


def run_text_overlap_scan(slide: Any, slide_num: int) -> list[Finding]:
    candidates = [
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and not getattr(shape, "has_table", False)
        and not getattr(shape, "has_chart", False)
        and shape_text(shape)
    ]
    findings: list[Finding] = []
    for index, first in enumerate(candidates):
        for second in candidates[index + 1:]:
            ratio = _intersection_ratio(first, second)
            if ratio >= 0.20:
                findings.append(Finding(
                    "warning", slide_num, "text_overlap",
                    "Independent text shapes overlap materially: "
                    f"'{getattr(first, 'name', 'unnamed')}' and '{getattr(second, 'name', 'unnamed')}' "
                    f"({ratio:.0%} of the smaller box)",
                ))
                if len(findings) >= 10:
                    return findings
    return findings


def run_qa(path: Path, facts_path: Path | None = None, briefs_path: Path | None = None) -> list[Finding]:
    prs = Presentation(str(path))
    facts = load_facts(facts_path)
    briefs_meta, briefs = load_briefs(briefs_path)
    findings: list[Finding] = []

    if not has_font_xml(path, "Arial"):
        findings.append(Finding("error", None, "font_xml", "Missing Arial font in slide XML"))
    if not has_font_xml(path, "Microsoft YaHei"):
        findings.append(Finding("error", None, "font_xml", "Missing Microsoft YaHei font in slide XML"))

    for index, slide in enumerate(prs.slides, start=1):
        texts = [shape_text(shape) for shape in slide.shapes if shape_text(shape)]
        all_text = "\n".join(texts)
        cover_or_divider = is_cover_or_divider(index, texts)

        if not cover_or_divider:
            title = first_top_text(slide)
            if not title:
                findings.append(Finding("error", index, "action_title", "Missing top action title"))
            else:
                if len(title) > 95:
                    findings.append(Finding("warning", index, "action_title", "Title may exceed two lines"))
                if looks_like_topic_label(title):
                    findings.append(Finding("warning", index, "action_title", "Title may be a topic label"))
            if len(re.findall(r"\d", all_text)) >= 8 and not re.search(r"(来源[:：]|Source:)\s*\S+", all_text):
                findings.append(Finding("warning", index, "source_line", "Numeric-heavy slide may need a source line"))
            findings.extend(run_qualitative_claim_scan(slide, index))

        if has_cjk(all_text):
            for term, Chinese in BANNED_TERMS.items():
                if re.search(r"\b" + re.escape(term) + r"\b", all_text.lower()):
                    findings.append(Finding("warning", index, "terminology", f"Consider Chinese term for {term}: {Chinese}"))

        for shape in slide.shapes:
            left, top = cm(shape.left), cm(shape.top)
            right, bottom = left + cm(shape.width), top + cm(shape.height)
            if left < -0.05 or top < -0.05 or right > SLIDE_W_CM + 0.05 or bottom > SLIDE_H_CM + 0.05:
                findings.append(Finding("error", index, "bounds", "Shape out of bounds"))

        findings.extend(run_text_overflow_scan(slide, index))
        findings.extend(run_text_overlap_scan(slide, index))

    findings.extend(run_fact_consistency(prs, facts))
    findings.extend(run_data_density(prs, facts, briefs_meta, briefs))
    return findings


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run automated consulting-deck PPTX QA")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--json", action="store_true")
    parser.add_argument("--facts", type=Path, default=None, help="Optional evidence.json fact table")
    parser.add_argument("--briefs", type=Path, default=None, help="Optional briefs.yaml for page roles and density thresholds")
    args = parser.parse_args(argv)
    findings = run_qa(args.pptx, args.facts, args.briefs)

    modules = [
        "fonts", "action-title", "source-line", "qualitative-claims",
        "terminology", "bounds", "text-overflow", "text-overlap",
    ]
    if args.facts is not None:
        modules.append("fact-consistency")
    if args.briefs is not None:
        modules.append("data-density")
    fact_count = len(load_facts(args.facts)) if args.facts is not None else 0
    brief_count = len(load_briefs(args.briefs)[1]) if args.briefs is not None else 0
    errors = sum(item.severity == "error" for item in findings)
    warnings = sum(item.severity == "warning" for item in findings)
    status = (
        f"QA passed: 0 findings; modules={','.join(modules)}; "
        f"facts={fact_count if args.facts is not None else 'not-supplied'}; "
        f"briefs={brief_count if args.briefs is not None else 'not-supplied'}"
        if not findings else
        f"QA completed: errors={errors}, warnings={warnings}; modules={','.join(modules)}; "
        f"facts={fact_count if args.facts is not None else 'not-supplied'}; "
        f"briefs={brief_count if args.briefs is not None else 'not-supplied'}"
    )

    if args.json:
        print(json.dumps([asdict(item) for item in findings], ensure_ascii=False, indent=2))
        print(status, file=sys.stderr)
    else:
        for finding in findings:
            where = f"slide {finding.slide}" if finding.slide is not None else "deck"
            print(f"[{finding.severity}] {where} {finding.check}: {finding.message}")
        print(status)
    return 1 if errors else 0


if __name__ == "__main__":
    raise SystemExit(main())
