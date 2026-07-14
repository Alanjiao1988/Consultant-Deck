"""Generate a one-page regulatory map demo and its semantic manifest."""
from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.jurisdiction_map import categorical_jurisdiction_map, write_exhibit_manifest

SLIDE_W = 33.87
SLIDE_H = 19.05


def _set_font(run, size, bold=False, color="333333"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = "Arial"
    from pptx.oxml.ns import qn
    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = r_pr.makeelement(qn("a:ea"), {})
        r_pr.append(ea)
    ea.set("typeface", "Microsoft YaHei")


def generate(
    output_path: Path | None = None,
    manifest_path: Path | None = None,
) -> tuple[Path, Path]:
    output_path = output_path or ROOT / "examples" / "demo_regulatory_map.pptx"
    manifest_path = manifest_path or ROOT / "examples" / "demo_regulatory_map.exhibits.json"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    kicker = slide.shapes.add_textbox(Cm(1.2), Cm(0.55), Cm(8.0), Cm(0.35))
    kicker.text_frame.margin_left = kicker.text_frame.margin_right = 0
    run = kicker.text_frame.paragraphs[0].add_run()
    run.text = "GLOBAL REGULATORY MAP"
    _set_font(run, 8.5, True, "595959")

    title = slide.shapes.add_textbox(Cm(1.2), Cm(1.05), Cm(31.0), Cm(1.25))
    title.text_frame.margin_left = title.text_frame.margin_right = 0
    title.text_frame.word_wrap = True
    run = title.text_frame.paragraphs[0].add_run()
    run.text = "Across 24 reviewed markets, strict localization remains concentrated in regulated workloads rather than the global default"
    _set_font(run, 19.0, True, "1F3864")

    categories = [
        {
            "id": "unified_rulebook",
            "label": "Unified rulebook",
            "description": "Bloc-level framework is the primary operating reference",
        },
        {
            "id": "federated_sectoral",
            "label": "Federated and sector-led",
            "description": "Federal, subnational and industry rules interact",
        },
        {
            "id": "sector_localization",
            "label": "Sector-specific localization",
            "description": "Government, finance, health or critical infrastructure",
        },
        {
            "id": "transfer_procedure",
            "label": "Transfer procedure",
            "description": "Mechanism, filing, consent or adequacy is decision-critical",
        },
    ]
    jurisdictions = [
        {
            "id": "US_CA",
            "label": "US / Canada",
            "category_id": "federated_sectoral",
            "note": "State, province and sector overlays",
            "evidence_ids": ["F102"],
            "entity_type": "custom_group",
            "members": ["US", "CA"],
        },
        {
            "id": "EU",
            "label": "European Union",
            "category_id": "unified_rulebook",
            "note": "Bloc rulebook with member-state implementation",
            "evidence_ids": ["F101"],
            "entity_type": "regulatory_bloc",
            "members": ["EU member states"],
        },
        {
            "id": "TR",
            "label": "Türkiye",
            "category_id": "transfer_procedure",
            "note": "Transfer mechanism and authority practice matter",
            "evidence_ids": ["F105"],
        },
        {
            "id": "UAE_SA",
            "label": "UAE / Saudi Arabia",
            "category_id": "sector_localization",
            "note": "Regulated sectors drive local deployment",
            "evidence_ids": ["F103"],
            "entity_type": "custom_group",
            "members": ["UAE", "SA"],
        },
        {
            "id": "BR",
            "label": "Brazil",
            "category_id": "transfer_procedure",
            "note": "Cross-border legal mechanism is the key gate",
            "evidence_ids": ["F104"],
        },
        {
            "id": "JP_KR",
            "label": "Japan / South Korea",
            "category_id": "transfer_procedure",
            "note": "Consent and transfer basis require review",
            "evidence_ids": ["F106"],
            "entity_type": "custom_group",
            "members": ["JP", "KR"],
        },
    ]

    manifest = categorical_jurisdiction_map(
        slide,
        jurisdictions,
        categories,
        1.2,
        3.0,
        31.2,
        14.5,
        classification_dimension="Dominant regulatory operating model for enterprise data localization and cross-border transfer",
        as_of="2026-07-14",
        caveat="Illustrative demo classifications; validate current entity, workload, sector and data-type requirements before use.",
        coverage={"reviewed": 24, "shown": 6, "selection_basis": "Representative operating patterns"},
        insight_annotations=[
            "Localization decisions are workload-specific, not country-wide defaults",
            "Grouping markets is useful for planning but not a substitute for legal review",
        ],
        decision_implication="Use the map to triage detailed jurisdiction and workload assessments.",
        legend_title="Dominant operating pattern",
        page=1,
    )

    source = slide.shapes.add_textbox(Cm(1.2), Cm(18.15), Cm(27.0), Cm(0.32))
    source.text_frame.margin_left = source.text_frame.margin_right = 0
    r = source.text_frame.paragraphs[0].add_run()
    r.text = "Source: illustrative evidence records F101–F106; Natural Earth public-domain landmass base"
    _set_font(r, 7.0, False, "595959")
    page_no = slide.shapes.add_textbox(Cm(31.0), Cm(18.05), Cm(1.4), Cm(0.35))
    page_no.text_frame.margin_left = page_no.text_frame.margin_right = 0
    page_no.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = page_no.text_frame.paragraphs[0].add_run()
    r.text = "01"
    _set_font(r, 8.5, True, "595959")

    prs.save(output_path)
    write_exhibit_manifest(manifest_path, [manifest])
    return output_path, manifest_path


if __name__ == "__main__":
    generated = generate()
    print("Generated:", *generated)
