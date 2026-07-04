"""Business-case exhibits for consulting decks."""
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from .consulting_shapes import PRIMARY, ACCENT, GRAY_TEXT, GRAY_2, GRAY_3, GRAY_LINE, GRAY_FILL, add_textbox, _solid, _line, native_chart, waterfall


def kpi_strip(slide, metrics, x, y, w, h=1.65):
    """metrics: [(label, value, note, sentiment)] where sentiment is good/bad/neutral."""
    gap = 0.25
    card_w = (w - gap * (len(metrics) - 1)) / len(metrics)
    for i, (label, value, note, sentiment) in enumerate(metrics):
        cx = x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(cx), Cm(y), Cm(card_w), Cm(h))
        card.adjustments[0] = 0.08
        _solid(card, GRAY_FILL)
        _line(card, GRAY_LINE, 0.6)
        color = PRIMARY if sentiment == "good" else (ACCENT if sentiment == "bad" else GRAY_TEXT)
        add_textbox(slide, label, cx + 0.25, y + 0.18, card_w - 0.5, 0.35, size=8.5, color=GRAY_2)
        add_textbox(slide, value, cx + 0.25, y + 0.52, card_w - 0.5, 0.6, size=17, bold=True, color=color)
        add_textbox(slide, note, cx + 0.25, y + 1.18, card_w - 0.5, 0.35, size=8, color=GRAY_3)


def business_case_summary(slide, metrics, yearly_cashflows, x, y, w, h, unit="$m"):
    kpi_strip(slide, metrics, x, y, w, 1.65)
    categories = [str(row[0]) for row in yearly_cashflows]
    net = [row[3] for row in yearly_cashflows]
    native_chart(slide, "bar", categories, [(f"Net cash flow ({unit})", net)], x, y + 2.15, w, h - 2.15,
                 show_legend=False, show_values=True, value_format="#,##0")


def tco_waterfall(slide, items, x, y, w, h, unit="m"):
    return waterfall(slide, items, x, y, w, h, unit=unit)


def sensitivity_table(slide, rows, columns, values, x, y, w, h, fmt="{:g}"):
    nr, nc = len(rows), len(columns)
    label_w = w * 0.24
    col_w = (w - label_w) / nc
    row_h = (h - 0.75) / nr
    add_textbox(slide, "Scenario", x, y, label_w, 0.7, 9, True, GRAY_2)
    for j, col in enumerate(columns):
        add_textbox(slide, str(col), x + label_w + j * col_w, y, col_w, 0.7, 9, True, GRAY_2, PP_ALIGN.CENTER)
    flat = [v for row in values for v in row]
    best = max(flat) if flat else None
    for i, row_name in enumerate(rows):
        ry = y + 0.75 + i * row_h
        add_textbox(slide, str(row_name), x, ry + 0.08, label_w - 0.1, row_h, 9)
        for j, value in enumerate(values[i]):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x + label_w + j * col_w), Cm(ry), Cm(col_w - 0.05), Cm(row_h - 0.05))
            _solid(cell, "D6DEEC" if value == best else "FFFFFF")
            _line(cell, GRAY_LINE, 0.5)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = fmt.format(value)
            from .consulting_shapes import set_font
            set_font(r, 9, value == best, PRIMARY if value == best else GRAY_TEXT)
