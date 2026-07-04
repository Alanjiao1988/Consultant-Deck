"""Generate a reusable consulting template PPTX."""
from pathlib import Path
import sys
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.consulting_layouts import SLIDE_W, SLIDE_H, cover_page, section_divider, add_action_title, add_footer, exec_summary_block
from scripts.consulting_shapes import add_textbox, native_chart
from scripts.architecture_helpers import layered_architecture


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def generate(output=ROOT / "assets" / "consulting_template.pptx"):
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_W)
    prs.slide_height = Cm(SLIDE_H)

    slide = blank(prs)
    cover_page(slide, "Deck title states the decision or recommendation", "Subtitle / client / date", "Client name", "YYYY-MM-DD")

    slide = blank(prs)
    add_action_title(slide, "Executive summary synthesizes the answer in three to four decision-relevant messages")
    exec_summary_block(slide, [("Context", "Write the most important situational insight here.", ["Supporting point one", "Supporting point two"]), ("Implication", "Write the business implication and decision requirement here.", ["Supporting point one", "Supporting point two"])])
    add_footer(slide, 2, "Source / calculation basis")

    slide = blank(prs)
    section_divider(slide, 1, "Section title is a conclusion, not a topic label", ["Context", "Analysis", "Recommendation"], 0)
    add_footer(slide, 3)

    slide = blank(prs)
    add_action_title(slide, "A single chart page uses the chart to prove the action title")
    native_chart(slide, "bar", ["A", "B", "C"], [("Series", [10, 18, 25])], 1.4, 4.0, 20, 8, show_legend=False)
    add_textbox(slide, "Key implication", 23.0, 4.0, 8.5, 0.5, size=11, bold=True)
    add_textbox(slide, "Use this right rail for so-what, caveat and decision implications.", 23.0, 4.7, 8.5, 2.0, size=10.5)
    add_footer(slide, 4, "Source / calculation basis")

    slide = blank(prs)
    add_action_title(slide, "Architecture pages separate experience, orchestration, data, model and control layers")
    layered_architecture(slide, [("Experience", ["Channel A", "Channel B", "User app"]), ("Orchestration", ["Workflow", "Tool calling", "Policy"]), ("Knowledge", ["Vector index", "Data store", "Catalog"]), ("Model", ["LLM", "Embedding", "Eval"]), ("Control", ["Audit", "DLP", "Monitoring"])], 1.4, 3.3, 30.5, 10.5)
    add_footer(slide, 5, "Source / calculation basis")

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(output)


if __name__ == "__main__":
    generate()
